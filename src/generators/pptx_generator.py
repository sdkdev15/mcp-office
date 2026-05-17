"""PowerPoint presentation generator using python-pptx with theme support."""

from __future__ import annotations

import io
import re
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from src.styles.themes import get_theme
from src.utils.colors import hex_to_rgbcolor_tuple
from src.utils.metadata import apply_metadata
from src.utils.logger import get_logger
from src.utils.validators import validate_slide_layout, ValidationError
from src.utils.image_handler import ImageHandler

log = get_logger("pptx_generator")
image_handler = ImageHandler()


def hex_to_rgbcolor(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor for python-pptx."""
    r, g, b = hex_to_rgbcolor_tuple(hex_color)
    return RGBColor(r, g, b)


class PPTXGenerator:
    """Generates PowerPoint presentations with rich slides, charts, and styling."""

    SLIDE_WIDTH_STANDARD_INCHES = 10.0
    SLIDE_HEIGHT_STANDARD_INCHES = 7.5
    SLIDE_WIDTH_WIDESCREEN_INCHES = 13.33
    SLIDE_HEIGHT_WIDESCREEN_INCHES = 7.5

    def __init__(self, theme_name: str = "corporate"):
        self.theme_name = theme_name
        self.theme = get_theme(theme_name)
        self.pres: Optional[Presentation] = None

    def create_presentation(
        self,
        title: str = "Presentation",
        slide_size: str = "widescreen",
        metadata: Optional[dict] = None,
        slides: Optional[list[dict]] = None,
    ) -> bytes:
        """Create a new PowerPoint presentation.

        Args:
            title: Presentation title.
            slide_size: Slide size (widescreen for 16:9, standard for 4:3).
            metadata: Optional document metadata.
            slides: Optional list of slides to add.

        Returns:
            Presentation content as bytes.
        """
        self.pres = Presentation()

        # Set slide size
        if slide_size.lower() == "standard":
            self.pres.slide_width = Inches(self.SLIDE_WIDTH_STANDARD_INCHES)
            self.pres.slide_height = Inches(self.SLIDE_HEIGHT_STANDARD_INCHES)
        else:  # widescreen
            self.pres.slide_width = Inches(self.SLIDE_WIDTH_WIDESCREEN_INCHES)
            self.pres.slide_height = Inches(self.SLIDE_HEIGHT_WIDESCREEN_INCHES)

        # Apply metadata
        if metadata:
            self._apply_metadata(metadata)

        # Add title slide
        self.add_slide("title", title=title)

        # Add content slides
        if slides:
            for slide_data in slides:
                self.add_slide(
                    layout=slide_data.get("layout", "title_and_content"),
                    title=slide_data.get("title"),
                    content=slide_data.get("content"),
                    bullets=slide_data.get("bullets"),
                    images=slide_data.get("images"),
                )

        return self._save_presentation()

    def add_slide(
        self,
        layout: str = "title_and_content",
        title: Optional[str] = None,
        content: Optional[str] = None,
        bullets: Optional[list[str]] = None,
        images: Optional[list[dict]] = None,
    ) -> int:
        """Add a slide to the presentation.

        Args:
            layout: Slide layout type.
            title: Slide title.
            content: Body content text.
            bullets: List of bullet points.
            images: List of images to add.

        Returns:
            Index of the added slide.
        """
        layout = validate_slide_layout(layout)
        slide_layout_map = {
            "title": 0,
            "title_and_content": 1,
            "blank": 6,
            "two_content": 2,
            "comparison": 5,
            "title_only": 5,
            "section_header": 11,
        }

        layout_idx = slide_layout_map.get(layout, 1)

        try:
            slide_layout = self.pres.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = self.pres.slide_layouts[1]

        slide = self.pres.slides.add_slide(slide_layout)
        slide_idx = len(self.pres.slides) - 1

        # Premium graphics are deferred to end of slide construction

        # Set title
        if title:
            try:
                if getattr(slide.shapes, 'title', None):
                    slide.shapes.title.text = title
                    self._style_title(slide.shapes.title)
            except Exception as e:
                log.warning(f"Failed to set title for slide '{title}': {e}")

        # Add content
        if content or bullets:
            try:
                body_shape = None
                # Find body placeholder by index (idx=1 is the standard body)
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        body_shape = ph
                        break
                if body_shape is None and len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]

                if body_shape is None:
                    raise ValueError("No body placeholder found")

                tf = body_shape.text_frame
                from pptx.enum.text import MSO_AUTO_SIZE
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.word_wrap = True

                # Clear existing paragraphs (keep first, remove rest, then clear first)
                for p in list(tf.paragraphs)[1:]:
                    tf._txBody.remove(p._p)
                # Clear the mandatory first paragraph's text
                if tf.paragraphs:
                    tf.paragraphs[0].clear()

                first_para = True
                if content:
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.text = content

                if bullets:
                    for bullet in bullets:
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        # Set bullet formatting via XML
                        pPr = p._p.get_or_add_pPr()
                        buChar = OxmlElement('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
            except Exception as e:
                log.warning(f"Failed to format body text: {e}")
                # Fallback: add a text box
                if content:
                    self.add_text_box(slide_idx, content, left=1, top=2)
                elif bullets:
                    self.add_text_box(slide_idx, "\n".join(bullets), left=1, top=2)
        self._draw_wow_graphics(slide)
        
        # Add images
        if images:
            for img in images:
                try:
                    src = img.get("source")
                    if src:
                        cached = image_handler.process_image(src)
                        self.add_image(
                            slide_idx, 
                            cached, 
                            left=img.get("left", 1.0), 
                            top=img.get("top", 1.0), 
                            width=img.get("width"), 
                            height=img.get("height")
                        )
                except Exception as e:
                    log.warning(f"Failed to add image to slide: {e}")

        return slide_idx

    def add_text_box(
        self,
        slide_index: int,
        text: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 5.0,
        height: float = 1.0,
        font_size: Optional[int] = None,
        bold: bool = False,
        color: Optional[str] = None,
    ) -> None:
        """Add a text box to a slide."""
        slide = self.pres.slides[slide_index]
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text

        run = p.runs[0] if p.runs else p.add_run()
        run.bold = bold
        if color:
            run.font.color.rgb = hex_to_rgbcolor(color)

    def add_image(
        self,
        slide_index: int,
        image_path: str,
        left: float = 1.0,
        top: float = 1.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> None:
        """Add an image to a slide."""
        slide = self.pres.slides[slide_index]
        slide.shapes.add_picture(
            image_path,
            Inches(left),
            Inches(top),
            width=Inches(width) if width else None,
            height=Inches(height) if height else None,
        )

    def add_table(
        self,
        slide_index: int,
        headers: list[str],
        rows: list[list[Any]],
    ) -> None:
        """Add a table to a slide."""
        slide = self.pres.slides[slide_index]
        num_rows = len(rows) + 1
        num_cols = len(headers)

        table_width = Inches(min(num_cols * 1.5, 10))
        table_height = Inches(min(num_rows * 0.4, 5))

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, Inches(1), Inches(1.5), table_width, table_height
        )
        table = table_shape.table
        colors = self.theme.colors

        # Set headers
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.bold = True
                    run.font.size = Pt(11)
                    run.font.color.rgb = hex_to_rgbcolor(colors.header_text)

        # Set data
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = hex_to_rgbcolor(colors.text)

    def add_chart(
        self,
        slide_index: int,
        chart_type: str,
        headers: list[str],
        rows: list[list[Any]],
        title: str = "Chart",
    ) -> None:
        """Add a chart to a slide."""
        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }

        chart_type_enum = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
        slide = self.pres.slides[slide_index]

        all_data = [headers] + [[str(cell) for cell in row] for row in rows]
        chart_data = self._create_chart_data(all_data)

        chart_frame = slide.shapes.add_chart(
            chart_type_enum, Inches(1), Inches(1), Inches(9), Inches(5), chart_data
        )
        chart = chart_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    def set_background(self, slide_index: int, color: Optional[str] = None) -> None:
        """Set slide background color."""
        slide = self.pres.slides[slide_index]
        fill = slide.background.fill
        fill.solid()
        bg_color = color or self.theme.colors.background
        fill.fore_color.rgb = hex_to_rgbcolor(bg_color)

    @staticmethod
    def _strip_markdown(text: str) -> str:
        """Strip Markdown formatting from text.
        
        Removes: **bold**, *italic*, `inline code`, [link](url), ~~strikethrough~~,
        heading markers (##), horizontal rules (---), and blockquote markers (>).
        
        Args:
            text: Text with Markdown formatting.
        
        Returns:
            Clean text without Markdown markers.
        """
        # Remove heading markers (###, ##, #) at start of line
        result = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
        
        # Remove horizontal rules
        result = re.sub(r'^[-*_]{3,}\s*$', '', result, flags=re.MULTILINE)
        
        # Remove blockquote markers
        result = re.sub(r'^>\s+', '', result, flags=re.MULTILINE)
        
        # Remove strikethrough ~~text~~
        result = re.sub(r'~~(.*?)~~', r'\1', result)
        
        # Remove **bold** or __bold__
        result = re.sub(r'\*\*(.+?)\*\*', r'\1', result, flags=re.DOTALL)
        result = re.sub(r'__(.+?)__', r'\1', result, flags=re.DOTALL)
        
        # Remove *italic* or _italic_
        result = re.sub(r'\*(.+?)\*', r'\1', result, flags=re.DOTALL)
        result = re.sub(r'(?<!\w)_(.+?)_(?!\w)', r'\1', result, flags=re.DOTALL)
        
        # Remove `inline code`
        result = re.sub(r'`([^`]+)`', r'\1', result)
        
        # Remove [link](url) -> link
        result = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', result)
        
        # Remove image syntax ![alt](url)
        result = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', result)
        
        return result.strip()

    def create_from_prompt(self, prompt: str, title: str = "Presentation") -> bytes:
        """Create a presentation from a natural language prompt.

        Parses the prompt to extract slides from headings, "Heading - content" lines,
        Markdown headings (##), and tables.
        """
        log.info(f"Creating presentation from prompt: {prompt[:100]}...")

        self.pres = Presentation()
        self.pres.slide_width = Inches(self.SLIDE_WIDTH_WIDESCREEN_INCHES)
        self.pres.slide_height = Inches(self.SLIDE_HEIGHT_WIDESCREEN_INCHES)

        # Extract title from prompt if generic
        doc_title = title if title != "Presentation" else self._extract_title(prompt)
        self.add_slide("title", title=doc_title)

        # Parse prompt into slides
        self._parse_prompt_to_slides(prompt)

        return self._save_presentation()

    def _extract_title(self, text: str) -> str:
        """Extract presentation title from prompt text."""
        lines = text.strip().split("\n")
        for line in lines:
            line = line.strip()
            # Skip instruction lines
            if re.match(r"^(Buat|Format|Gunakan|Create|Generate)", line, re.IGNORECASE):
                continue
            # Strip Markdown heading markers
            line = re.sub(r'^#{1,6}\s+', '', line)
            # Strip Markdown formatting
            line = self._strip_markdown(line)
            if line and len(line) > 10:
                return line.rstrip(":").strip()
        return "Presentation"

    def _parse_prompt_to_slides(self, text: str) -> None:
        """Parse prompt text into slides.
        
        Supports:
        - Markdown headings (##, ###) as slide titles
        - "Heading - content" pattern
        - ALL CAPS headings
        - Bullet points (appended to last slide)
        - Tables (pipe-separated) as table slides
        - Regular paragraphs as content slides
        """
        lines = text.strip().split("\n")
        pending_bullets = []
        table_buffer = []

        def flush_bullets():
            """Add pending bullets to the last slide as content."""
            nonlocal pending_bullets
            if pending_bullets and len(self.pres.slides) > 0:
                # Append bullets to last slide's content
                last_slide = self.pres.slides[-1]
                bullets_text = "\n".join(["- " + b for b in pending_bullets])
                # Add as a new content paragraph if there's room
                try:
                    for shape in last_slide.shapes:
                        if shape.has_text_frame:
                            tf = shape.text_frame
                            p = tf.add_paragraph()
                            p.text = bullets_text
                            for run in p.runs:
                                self._style_run(run)
                            break
                except Exception:
                    pass
            pending_bullets = []

        def flush_table():
            """Create a slide from buffered table data."""
            nonlocal table_buffer
            if len(table_buffer) >= 2:
                headers = table_buffer[0]
                rows = table_buffer[1:]
                # Create a slide with the table
                slide_layout_map = {
                    "title_and_content": 1,
                    "blank": 6,
                }
                try:
                    slide_layout = self.pres.slide_layouts[1]
                except IndexError:
                    slide_layout = self.pres.slide_layouts[0]
                slide = self.pres.slides.add_slide(slide_layout)
                if headers and getattr(slide.shapes, 'title', None):
                    slide.shapes.title.text = headers[0] if len(headers) > 1 else "Table"
                    self._style_title(slide.shapes.title)
                # Add table to slide
                try:
                    self.add_table(len(self.pres.slides) - 1, headers, rows)
                except Exception as e:
                    log.warning(f"Failed to add table to slide: {e}")
                    # Fallback: add as text
                    self.add_slide(title="Table", content="\n".join([" | ".join(str(c) for c in row) for row in [headers] + rows[:5]]))
            table_buffer = []

        i = 0
        while i < len(lines):
            line = lines[i].strip()
            i += 1

            # Skip empty lines
            if not line:
                continue

            # Skip instruction lines
            if re.match(r"^(Buat|Format|Gunakan|Create|Generate|Silakan|Please)", line, re.IGNORECASE):
                continue

            # Skip Markdown table separator lines (|---|---|)
            if re.match(r'^\|[\s\-:|]+\|$', line):
                continue

            # Detect Markdown table rows
            if "|" in line and line.count("|") >= 2:
                cells = [self._strip_markdown(c.strip()) for c in line.split("|") if c.strip()]
                if cells:
                    table_buffer.append(cells)
                flush_bullets()
                continue

            # If we hit non-table content, flush table buffer
            if table_buffer and "|" not in line:
                flush_table()

            # Detect Markdown headings (##, ###, ####)
            md_heading = re.match(r'^#{1,6}\s+(.+)$', line)
            if md_heading:
                flush_bullets()
                flush_table()
                heading_text = self._strip_markdown(md_heading.group(1).strip())
                self.add_slide(title=heading_text)
                continue

            # Detect "Heading - content" pattern
            dash_match = re.match(r"^([A-Z][^-\n]{5,}?)\s+-\s+(.+)$", line)
            if dash_match:
                flush_bullets()
                flush_table()
                slide_title = self._strip_markdown(dash_match.group(1).strip())
                content_text = self._strip_markdown(dash_match.group(2).strip())

                # Try to split content into bullets (comma-separated)
                if "," in content_text and len(content_text) > 50:
                    parts = [self._strip_markdown(p.strip()) for p in content_text.split(",") if p.strip()]
                    if len(parts) > 1:
                        self.add_slide(title=slide_title, bullets=parts)
                        continue

                self.add_slide(title=slide_title, content=content_text)
                continue

            # ALL CAPS heading
            if line.isupper() and not line.startswith(("-", "*", "•")):
                flush_bullets()
                flush_table()
                self.add_slide(title=line.title())
                continue

            # Bullet points
            bullet_match = re.match(r"^[-*•]\s+(.*)", line)
            if bullet_match:
                bullet_text = self._strip_markdown(bullet_match.group(1).strip())
                if bullet_text:
                    pending_bullets.append(bullet_text)
                continue

            # Numbered items
            num_match = re.match(r"^\d+\.\s+(.*)", line)
            if num_match:
                bullet_text = self._strip_markdown(num_match.group(1).strip())
                if bullet_text:
                    pending_bullets.append(bullet_text)
                continue

            # Regular paragraph - treat as content slide if long enough
            if len(line) > 20:
                flush_bullets()
                flush_table()
                clean_line = self._strip_markdown(line)
                if clean_line:
                    self.add_slide(title="Content", content=clean_line)

        # Flush remaining buffers
        flush_bullets()
        flush_table()

    def create_from_template(
        self,
        template_path: str,
        slides: Optional[list[dict]] = None,
        metadata: Optional[dict] = None,
    ) -> bytes:
        """Create a presentation using an existing .pptx file as template.

        The template preserves all master slide designs, themes, fonts,
        and layouts. New slides are appended after existing content.

        Args:
            template_path: Path to the .pptx template file.
            slides: Optional list of slide dicts with title, content, bullets.
            metadata: Optional document metadata.

        Returns:
            Presentation content as bytes.
        """
        self.pres = Presentation(template_path)

        # Apply metadata
        if metadata:
            self._apply_metadata(metadata)

        # Add new slides
        if slides:
            for slide_data in slides:
                self.add_slide(
                    layout=slide_data.get("layout", "title_and_content"),
                    title=slide_data.get("title"),
                    content=slide_data.get("content"),
                    bullets=slide_data.get("bullets"),
                    images=slide_data.get("images"),
                )

        return self._save_presentation()

    def _style_title(self, shape) -> None:
        """Apply theme styling to a title shape."""
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            # If no runs exist, add one
            for run in paragraph.runs:
                pass

        pass

    def _draw_wow_graphics(self, slide) -> None:
        """Inject theme-specific vector graphics onto the slide."""
        layout_map = {
            "corporate": self._layout_corporate,
            "minimal": self._layout_minimal,
            "creative": self._layout_creative,
            "academic": self._layout_academic,
            "dark": self._layout_dark,
        }
        handler = layout_map.get(self.theme_name, self._layout_corporate)
        handler(slide)

    def _send_to_back(self, element, index=0):
        """Send a shape element to a specific z-order index.
        Note: In OOXML <p:spTree>, the first two children MUST be <p:nvGrpSpPr> and <p:grpSpPr>.
        Therefore, valid shape elements must be inserted starting at index 2 to prevent MS PowerPoint corruption.
        """
        element.getparent().insert(2 + index, element)

    def _style_title_white(self, slide):
        """Override title text to white."""
        from pptx.util import Inches
        if getattr(slide.shapes, 'title', None):
            try:
                t = slide.shapes.title
                t.top = Inches(0.2)
                t.left = Inches(0.5)
                t.width = Inches(12.33)
                t.height = Inches(0.8)
                if t.has_text_frame:
                    for p in t.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = hex_to_rgbcolor("FFFFFF")
            except Exception:
                pass

    def _style_title_color(self, slide, color: str, left=0.5, top=0.15, width=12.33, height=0.8):
        """Override title text to a specific color and position."""
        from pptx.util import Inches
        if getattr(slide.shapes, 'title', None):
            try:
                t = slide.shapes.title
                t.top = Inches(top)
                t.left = Inches(left)
                t.width = Inches(width)
                t.height = Inches(height)
                if t.has_text_frame:
                    for p in t.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = hex_to_rgbcolor(color)
            except Exception:
                pass

    # ── Corporate: Full-width header bar + accent line + footer ──
    def _layout_corporate(self, slide):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        c = self.theme.colors

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.2))
        bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgbcolor(c.primary); bg.line.fill.background()
        self._send_to_back(bg._element, 0)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(13.33), Inches(0.05))
        accent.fill.solid(); accent.fill.fore_color.rgb = hex_to_rgbcolor(c.accent); accent.line.fill.background()
        self._send_to_back(accent._element, 1)

        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.33), Inches(0.4))
        footer.fill.solid(); footer.fill.fore_color.rgb = hex_to_rgbcolor(c.table_alt_row); footer.line.fill.background()
        self._send_to_back(footer._element, 0)

        self._style_title_white(slide)

    # ── Minimal: Thin left accent bar + bottom rule ──
    def _layout_minimal(self, slide):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5))
        sidebar.fill.solid(); sidebar.fill.fore_color.rgb = hex_to_rgbcolor("000000"); sidebar.line.fill.background()
        self._send_to_back(sidebar._element, 0)

        bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.02))
        bottom.fill.solid(); bottom.fill.fore_color.rgb = hex_to_rgbcolor("E5E5E5"); bottom.line.fill.background()
        self._send_to_back(bottom._element, 0)

        self._style_title_color(slide, "111111", left=0.6, top=0.4)

    # ── Creative: Warm cream background + amber header + gold accents ──
    def _layout_creative(self, slide):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        c = self.theme.colors

        # Cream background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgbcolor(c.background); bg.line.fill.background()
        self._send_to_back(bg._element, 0)

        # Warm amber header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.4))
        header.fill.solid(); header.fill.fore_color.rgb = hex_to_rgbcolor(c.primary); header.line.fill.background()
        self._send_to_back(header._element, 1)

        # Gold accent line below header
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.4), Inches(13.33), Inches(0.06))
        strip.fill.solid(); strip.fill.fore_color.rgb = hex_to_rgbcolor(c.accent); strip.line.fill.background()
        self._send_to_back(strip._element, 2)

        # Thin gold bottom line
        bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.03))
        bottom.fill.solid(); bottom.fill.fore_color.rgb = hex_to_rgbcolor(c.accent); bottom.line.fill.background()
        self._send_to_back(bottom._element, 1)

        # Title: cream text on amber header
        self._style_title_color(slide, c.header_text)

    # ── Academic: Top double rule + classic positioning ──
    def _layout_academic(self, slide):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        c = self.theme.colors

        # Thick top rule
        top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.06))
        top_rule.fill.solid(); top_rule.fill.fore_color.rgb = hex_to_rgbcolor(c.primary); top_rule.line.fill.background()
        self._send_to_back(top_rule._element, 0)

        # Thin rule below
        thin_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(12.33), Inches(0.02))
        thin_rule.fill.solid(); thin_rule.fill.fore_color.rgb = hex_to_rgbcolor(c.primary); thin_rule.line.fill.background()
        self._send_to_back(thin_rule._element, 1)

        # Bottom matching rule
        bot_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.04))
        bot_rule.fill.solid(); bot_rule.fill.fore_color.rgb = hex_to_rgbcolor(c.primary); bot_rule.line.fill.background()
        self._send_to_back(bot_rule._element, 0)

        self._style_title_color(slide, c.primary, left=0.5, top=0.55)

    # ── Dark: Full-bleed dark background + accent stripe ──
    def _layout_dark(self, slide):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        c = self.theme.colors

        # Full slide dark background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgbcolor(c.header_bg); bg.line.fill.background()
        self._send_to_back(bg._element, 0)

        # Accent stripe at top
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.08))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = hex_to_rgbcolor(c.accent); stripe.line.fill.background()
        self._send_to_back(stripe._element, 1)

        # Bottom accent line
        bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), Inches(13.33), Inches(0.08))
        bottom.fill.solid(); bottom.fill.fore_color.rgb = hex_to_rgbcolor(c.accent); bottom.line.fill.background()
        self._send_to_back(bottom._element, 1)

        self._style_title_white(slide)

        # Make ALL text on the slide white for readability
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = hex_to_rgbcolor("F1F5F9")

    def _apply_metadata(self, metadata: dict) -> None:
        """Apply presentation metadata."""
        apply_metadata(self.pres.core_properties, metadata)

    def _create_chart_data(self, data: list[list]) -> object:
        """Create chart data from 2D array."""
        from pptx.chart.data import CategoryChartData

        chart_data = CategoryChartData()
        chart_data.categories = [str(item) for item in data[0][1:]] if len(data) > 1 else []

        for row in data[1:]:
            chart_data.add_series(str(row[0]), [float(x) if x else 0 for x in row[1:]])

        return chart_data

    def _save_presentation(self) -> bytes:
        """Save presentation to bytes."""
        buffer = io.BytesIO()
        self.pres.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
