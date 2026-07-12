"""Visual slide builder -- creates premium slides with auto-icons and visual elements.

Works alongside PPTXGenerator to build cover slides, agenda slides, exec-summary slides,
timeline slides, and flow diagram slides with premium layouts.
"""

from __future__ import annotations

from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from src.visual.gradient_engine import GradientEngine, GradientDef
from src.visual.shadow_engine import ShadowEngine, ShadowDef
from src.utils.auto_size import AutoSizeEngine

from src.styles.themes import get_theme, Theme
from src.utils.colors import hex_to_rgbcolor_tuple
from src.utils.logger import get_logger
from src.utils.visual_elements import VisualElements
from src.utils.icon_library import find_icon_by_keyword, get_icon
from src.utils.icon_renderer import IconRenderer

log = get_logger("visual_slide_builder")

# ── Constants ──
SLIDE_W = 13.33
SLIDE_H = 7.5
CARD_CORNER = 12  # Pt


class VisualSlideBuilder:
    """Build premium slides with auto-icons and visual elements.

    Usage:
        builder = VisualSlideBuilder(theme_name="corporate")
        builder.build_cover_slide(pres, title="...", subtitle="...")
    """

    def __init__(self, theme_name: str = "corporate"):
        self.theme_name = theme_name
        self.theme = get_theme(theme_name)
        self.visual = VisualElements()
        self.icon_renderer = IconRenderer()

    # ── Color helpers ──

    def _dark_bg(self) -> str:
        return self.theme.colors.header_bg

    def _accent(self) -> str:
        return self.theme.colors.accent

    def _text_on_dark(self) -> str:
        return self.theme.colors.header_text

    def _text_light_on_dark(self) -> str:
        return "#94A3B8"  # slate-400

    def _card_bg_dark(self) -> str:
        return "#1E293B"  # slightly lighter than bg

    def _text_on_light(self) -> str:
        return self.theme.colors.text

    def _text_light_on_light(self) -> str:
        return self.theme.colors.text_light

    # ── Shape helpers ──

    def _add_rect(
        self, slide, x: float, y: float, w: float, h: float,
        fill: str, send_back: int = 0,
    ):
        """Add a filled rectangle."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(fill))
        shape.line.fill.background()
        if send_back >= 0:
            self._send_to_back(shape._element, send_back)
        return shape

    def _add_rounded_rect(
        self, slide, x: float, y: float, w: float, h: float,
        fill: str, border: Optional[str] = None, border_width: int = 1,
    ):
        """Add a filled rounded rectangle."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(fill))
        if border:
            shape.line.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(border))
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()
        return shape

    def _add_text(
        self, slide, text: str,
        x: float, y: float, w: float, h: float,
        font_size: int = 14,
        bold: bool = False,
        color: str = "#FFFFFF",
        align: str = PP_ALIGN.LEFT,
        font_name: Optional[str] = None,
        italic: bool = False,
        gradient: Optional[dict] = None,
        shadow: Optional[dict] = None,
    ):
        """Add a text box with styling."""
        box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(color))
            if font_name:
                run.font.name = font_name
        
        # Apply gradient to text if specified
        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_text_frame(tf, grad)
        
        # Apply shadow to text if specified
        if shadow:
            shd = ShadowDef(**shadow)
            for run in p.runs:
                ShadowEngine.apply_to_text_run(run, shd)
        return box

    def _add_multiline_text(
        self, slide,
        lines: list[tuple[str, int, bool, str]],  # (text, size, bold, color)
        x: float, y: float, w: float, h: float,
        align: str = PP_ALIGN.LEFT,
        line_spacing: float = 1.3,
    ):
        """Add a text box with multiple styled lines."""
        box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        tf = box.text_frame
        tf.word_wrap = True

        for idx, (text, size, bold, color) in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.alignment = align
            p.space_after = Pt(int(size * line_spacing))
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(color))
                run.font.name = self.theme.fonts.body
        return box

    def _add_icon_circle(
        self, slide, icon_name: str, number: Optional[int] = None,
        x: float = 0, y: float = 0, size: float = 1.0,
        icon_color: str = "#F59E0B", bg_color: str = "#0F172A",
    ):
        """Add an icon inside a colored circle."""
        if number is not None:
            icon_path = self.icon_renderer.render_numbered_circle(
                number=number,
                size=int(size * 96),
                bg_color=icon_color,
                text_color=bg_color,
            )
        elif icon_name:
            icon_path = self.icon_renderer.render(
                icon_name=icon_name,
                size=int(size * 96),
                color=icon_color,
                bg_color=bg_color,
                bg_shape="circle",
            )
        else:
            return
        try:
            slide.shapes.add_picture(
                icon_path,
                Inches(x), Inches(y),
                Inches(size), Inches(size),
            )
        except Exception as e:
            log.warning(f"Failed to add icon: {e}")

    def _add_arrow(
        self, slide, from_x: float, from_y: float, to_x: float, to_y: float,
        color: str = "#F59E0B",
    ):
        """Add a simple arrow connector between two points."""
        arrow = slide.shapes.add_connector(
            1,  # straight connector
            Inches(from_x), Inches(from_y),
            Inches(to_x), Inches(to_y),
        )
        arrow.line.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(color))
        arrow.line.width = Pt(2)

    def _send_to_back(self, element, index: int = 0) -> None:
        """Send a shape element to the back (z-order)."""
        element.getparent().insert(2 + index, element)

    # ── Cover Slide ──

    def build_cover_slide(
        self,
        pres: Presentation,
        title: str,
        subtitle: str = "",
        icon_name: Optional[str] = None,
        extra_line: str = "",
    ) -> None:
        """Build a cover slide: dark bg + auto-detected icon + title + subtitle."""
        c = self.theme.colors

        if not icon_name:
            icon_name = find_icon_by_keyword(title) or find_icon_by_keyword(subtitle)

        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        # Full dark background
        self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        # Accent stripe at bottom
        self._add_rect(slide, 0, SLIDE_H - 0.1, SLIDE_W, 0.1, c.accent, send_back=1)

        # Decorative left accent bar
        self._add_rect(slide, 0.6, 2.5, 0.08, 2.5, c.accent, send_back=2)

        # Icon (centered near title)
        if icon_name:
            icon_x = 9.5
            icon_y = 1.5
            icon_size = 1.5
            # Circle bg for icon
            self._add_rounded_rect(
                slide, icon_x - 0.1, icon_y - 0.1, icon_size + 0.2, icon_size + 0.2,
                fill=c.header_bg,  # transparent-looking
            )
            self._add_icon_circle(
                slide, icon_name,
                x=icon_x, y=icon_y,
                size=icon_size,
                icon_color=c.accent,
                bg_color=c.header_bg,
            )

        # Title — large, bold
        self._add_text(
            slide, title,
            x=1.0, y=2.2, w=8.5, h=2.5,
            font_size=42, bold=True,
            color="#FFFFFF",
            font_name=self.theme.fonts.heading,
        )

        # Subtitle
        if subtitle:
            self._add_text(
                slide, subtitle,
                x=1.0, y=4.5, w=10, h=1.0,
                font_size=18, italic=True,
                color=self._text_light_on_dark(),
                font_name=self.theme.fonts.body,
            )

        # Extra line
        if extra_line:
            self._add_rect(slide, 1.0, 5.6, 3, 0.04, c.accent)
            self._add_text(
                slide, extra_line,
                x=1.0, y=5.8, w=10, h=0.5,
                font_size=12,
                color=self._text_light_on_dark(),
            )

    # ── Agenda Slide (Premium) ──

    def build_agenda_slide(
        self,
        pres: Presentation,
        title: str = "Agenda",
        items: Optional[list[dict]] = None,
    ) -> None:
        """Build an agenda slide with large numbered cards in a grid."""
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        # Header bar (full-width)
        self._add_rect(slide, 0, 0, SLIDE_W, 1.15, c.header_bg, send_back=3)
        self._add_rect(slide, 0, 1.15, SLIDE_W, 0.06, c.accent, send_back=2)

        # Title on header
        self._add_text(
            slide, title,
            x=0.8, y=0.25, w=11, h=0.7,
            font_size=36, bold=True,
            color="#FFFFFF",
            font_name=self.theme.fonts.heading,
        )

        if not items:
            return

        # Grid: 2 columns, items // 2 rows
        num = len(items)
        cols = 2
        rows = (num + cols - 1) // cols
        card_w = 5.4
        card_h = (SLIDE_H - 2.0) / rows - 0.2  # fill available height
        gap_x = 0.6
        start_x = 0.7
        start_y = 1.5

        for i, item in enumerate(items):
            col = i % cols
            row = i // cols
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + 0.2)

            # Card background
            self._add_rounded_rect(
                slide, x, y, card_w, card_h,
                fill=c.table_alt_row,
                border=c.border,
            )

            # Numbered circle — large
            num_size = 0.7
            num_x = x + 0.3
            num_y = y + 0.3
            self._add_icon_circle(
                slide, "",
                number=i + 1,
                x=num_x, y=num_y,
                size=num_size,
                icon_color=c.accent,
                bg_color=c.header_bg,
            )

            # Item title — bold
            item_title = str(item.get("title", f"Item {i + 1}"))
            self._add_text(
                slide, item_title,
                x=x + 1.2, y=y + 0.25,
                w=card_w - 1.5, h=0.5,
                font_size=16, bold=True,
                color=c.text,
                font_name=self.theme.fonts.heading,
            )

            # Item subtitle
            subtitle = item.get("subtitle", "")
            if subtitle:
                self._add_text(
                    slide, str(subtitle),
                    x=x + 1.2, y=y + 0.7,
                    w=card_w - 1.5, h=0.6,
                    font_size=12,
                    color=self._text_light_on_light(),
                )

    # ── Exec Summary Slide (Premium) ──

    def build_exec_summary_slide(
        self,
        pres: Presentation,
        title: str = "Ringkasan Eksekutif",
        stat_boxes: Optional[list[dict]] = None,
        body_text: str = "",
    ) -> None:
        """Build an exec summary slide with full-height stat column cards.

        Layout mirrors 'Mirrorbox' style: dark bg, large rounded cards
        that fill most of the slide height, accent-colored headers.
        """
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        # Full dark background
        self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=3)

        # Accent top stripe
        self._add_rect(slide, 0, 0, SLIDE_W, 0.06, c.accent, send_back=2)

        # Title
        self._add_text(
            slide, title,
            x=0.8, y=0.3, w=11, h=0.8,
            font_size=38, bold=True,
            color="#FFFFFF",
            font_name=self.theme.fonts.heading,
        )

        if not stat_boxes:
            return

        num_boxes = len(stat_boxes)
        # Layout: columns of equal width, full-height cards
        gap = 0.3
        total_gap = (num_boxes - 1) * gap + 1.6  # left + right margin
        card_w = (SLIDE_W - total_gap) / num_boxes
        card_h = SLIDE_H - 3.2  # fill most of remaining height
        start_x = 0.8
        start_y = 1.3

        # Accent colors for variety
        accent_colors = [c.accent, c.success, c.error, c.primary]

        for i, box_data in enumerate(stat_boxes):
            x = start_x + i * (card_w + gap)
            y = start_y

            box_color = box_data.get("bg_color", accent_colors[i % len(accent_colors)])

            # Full-height card
            self._add_rounded_rect(
                slide, x, y, card_w, card_h,
                fill=self._card_bg_dark(),
                border=box_color,
                border_width=2,
            )

            # Accent bar at top of card
            self._add_rect(
                slide, x + 0.1, y + 0.05, card_w - 0.2, 0.06,
                box_color,
            )

            # Number — large and bold, centered
            number = str(box_data.get("number", "0"))
            self._add_text(
                slide, number,
                x=x, y=y + 0.3,
                w=card_w, h=0.8,
                font_size=44, bold=True,
                color="#FFFFFF",
                align=PP_ALIGN.CENTER,
                font_name=self.theme.fonts.heading,
            )

            # Label — below number
            label = str(box_data.get("label", ""))
            if label:
                self._add_text(
                    slide, label,
                    x=x, y=y + 1.2,
                    w=card_w, h=0.5,
                    font_size=16, bold=True,
                    color=box_color,
                    align=PP_ALIGN.CENTER,
                    font_name=self.theme.fonts.heading,
                )

            # Sub-label
            sub_label = box_data.get("sub_label", "")
            if sub_label:
                self._add_text(
                    slide, str(sub_label),
                    x=x, y=y + 1.75,
                    w=card_w, h=0.5,
                    font_size=12,
                    color=self._text_light_on_dark(),
                    align=PP_ALIGN.CENTER,
                )

            # Icon at bottom of card
            icon_name = box_data.get("icon")
            if not icon_name:
                icon_name = find_icon_by_keyword(label)
            if icon_name:
                icon_size = 0.7
                icon_x = x + (card_w - icon_size) / 2
                icon_y = y + card_h - icon_size - 0.3
                self._add_icon_circle(
                    slide, icon_name,
                    x=icon_x, y=icon_y,
                    size=icon_size,
                    icon_color=box_color,
                    bg_color=self._card_bg_dark(),
                )

        # Body text at bottom
        if body_text:
            self._add_text(
                slide, body_text,
                x=1.0, y=SLIDE_H - 1.2,
                w=SLIDE_W - 2, h=0.8,
                font_size=13,
                color=self._text_light_on_dark(),
            )

    # ── Timeline Slide (Premium) ──

    def build_timeline_slide(
        self,
        pres: Presentation,
        title: str = "Timeline",
        events: Optional[list[dict]] = None,
        horizontal: bool = True,
    ) -> None:
        """Build a timeline slide with dark bg and full-width card layout."""
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        # Full dark background
        self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=4)

        # Accent top stripe
        self._add_rect(slide, 0, 0, SLIDE_W, 0.06, c.accent, send_back=3)

        # Title
        self._add_text(
            slide, title,
            x=0.8, y=0.3, w=11, h=0.8,
            font_size=38, bold=True,
            color="#FFFFFF",
            font_name=self.theme.fonts.heading,
        )

        if not events:
            return

        if horizontal:
            self._build_premium_horizontal_timeline(slide, events, c)
        else:
            self._build_vertical_timeline(slide, events, c)

    def _build_premium_horizontal_timeline(
        self, slide, events: list[dict], c: "ThemeColors",
    ) -> None:
        """Premium horizontal timeline: full-height cards with numbers."""
        num = len(events)
        if num == 0:
            return

        gap = 0.25
        margin = 0.7
        arrows_w = (num - 1) * 0.35 if num > 1 else 0
        total_content_w = SLIDE_W - margin * 2 - arrows_w
        card_w = (total_content_w - (num - 1) * gap) / num
        card_h = SLIDE_H - 3.0
        start_x = margin
        start_y = 1.4

        for i, event in enumerate(events):
            x = start_x + i * (card_w + gap + (0.35 if i < num - 1 else 0))
            y = start_y

            # Card background
            self._add_rounded_rect(
                slide, x, y, card_w, card_h,
                fill=self._card_bg_dark(),
            )

            # Numbered circle — large, top center
            num_size = 0.8
            num_x = x + (card_w - num_size) / 2
            num_y = y + 0.3
            self._add_icon_circle(
                slide, "",
                number=i + 1,
                x=num_x, y=num_y,
                size=num_size,
                icon_color=c.accent,
                bg_color=self._card_bg_dark(),
            )

            # Title
            title_text = str(event.get("title", ""))
            self._add_text(
                slide, title_text,
                x=x, y=y + 1.3,
                w=card_w, h=0.7,
                font_size=16, bold=True,
                color="#FFFFFF",
                align=PP_ALIGN.CENTER,
                font_name=self.theme.fonts.heading,
            )

            # Date
            date_text = event.get("date", "")
            if date_text:
                self._add_text(
                    slide, str(date_text),
                    x=x, y=y + 2.0,
                    w=card_w, h=0.4,
                    font_size=12,
                    color=c.accent,
                    align=PP_ALIGN.CENTER,
                )

            # Description
            desc_text = event.get("description", "")
            if desc_text:
                self._add_text(
                    slide, str(desc_text),
                    x=x + 0.2, y=y + 2.5,
                    w=card_w - 0.4, h=card_h - 3.0,
                    font_size=11,
                    color=self._text_light_on_dark(),
                    align=PP_ALIGN.CENTER,
                )

            # Arrow to next card
            if i < num - 1:
                arrow_x = x + card_w
                arrow_y = y + card_h / 2
                self._add_arrow(
                    slide, arrow_x, arrow_y,
                    arrow_x + 0.35, arrow_y,
                    color=c.accent,
                )

    # ── Flow Slide (Premium) ──

    def build_flow_slide(
        self,
        pres: Presentation,
        title: str = "Architecture",
        nodes: Optional[list[dict]] = None,
        connections: Optional[list[tuple[int, int]]] = None,
        columns: Optional[int] = None,
    ) -> None:
        """Build a flow diagram slide with configurable column layout.

        Args:
            pres: Presentation object.
            title: Slide title.
            nodes: List of dicts with keys: label, sub_label (optional), icon (optional).
            connections: List of (from_index, to_index) tuples.
            columns: Number of columns. Auto-calculated if None.
        """
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        # Header bar
        self._add_rect(slide, 0, 0, SLIDE_W, 1.15, c.header_bg, send_back=4)
        self._add_rect(slide, 0, 1.15, SLIDE_W, 0.06, c.accent, send_back=3)

        # Title
        self._add_text(
            slide, title,
            x=0.8, y=0.25, w=11, h=0.7,
            font_size=36, bold=True,
            color="#FFFFFF",
            font_name=self.theme.fonts.heading,
        )

        if not nodes:
            return

        num = len(nodes)
        if columns is None:
            columns = min(num, 5)
        rows = (num + columns - 1) // columns

        # Calculate card sizes to fill slide
        margin_x = 0.6
        margin_y = 1.5
        gap_x = 0.3
        gap_y = 0.5
        arrow_w = 0.4

        avail_w = SLIDE_W - margin_x * 2 - (columns - 1) * gap_x - (columns - 1) * arrow_w
        card_w = avail_w / columns
        avail_h = SLIDE_H - margin_y - 1.0  # leave room for bottom
        if rows > 1:
            avail_h -= (rows - 1) * gap_y
        card_h = min(avail_h, 4.5)  # cap height
        row_y_start = margin_y + (avail_h - card_h) / 2

        positions = []
        for i, node_data in enumerate(nodes):
            col = i % columns
            row = i // columns
            x = margin_x + col * (card_w + gap_x + arrow_w)
            y = row_y_start + row * (card_h + gap_y)
            positions.append((x, y))

            # Card background
            self._add_rounded_rect(
                slide, x, y, card_w, card_h,
                fill=c.table_alt_row,
                border=c.border,
            )

            # Numbered circle
            num_size = 0.6
            num_x = x + (card_w - num_size) / 2
            num_y = y + 0.2
            self._add_icon_circle(
                slide, "",
                number=i + 1,
                x=num_x, y=num_y,
                size=num_size,
                icon_color=c.accent,
                bg_color=c.header_bg,
            )

            # Title — bold
            label = str(node_data.get("label", f"Step {i + 1}"))
            self._add_text(
                slide, label,
                x=x, y=y + 1.0,
                w=card_w, h=0.6,
                font_size=14, bold=True,
                color=c.text,
                align=PP_ALIGN.CENTER,
                font_name=self.theme.fonts.heading,
            )

            # Sub-label / description
            sub_label = node_data.get("sub_label", "")
            if sub_label:
                self._add_text(
                    slide, str(sub_label),
                    x=x + 0.2, y=y + 1.6,
                    w=card_w - 0.4, h=card_h - 2.0,
                    font_size=11,
                    color=self._text_light_on_light(),
                    align=PP_ALIGN.CENTER,
                )

            # Arrow to next in row
            if i < num - 1:
                next_col = (i + 1) % columns
                if next_col == 0 or i + 1 >= num:
                    break
                # Only draw arrow if next item is in the next column (same row)
                next_row = (i + 1) // columns
                if next_row == row:
                    arrow_x = x + card_w
                    arrow_y = y + card_h / 2
                    self._add_arrow(
                        slide, arrow_x, arrow_y,
                        arrow_x + arrow_w, arrow_y,
                        color=c.accent,
                    )

        # Draw explicit connections
        if connections:
            for from_idx, to_idx in connections:
                if 0 <= from_idx < len(positions) and 0 <= to_idx < len(positions):
                    fx, fy = positions[from_idx]
                    tx, ty = positions[to_idx]
                    self._add_arrow(
                        slide,
                        fx + card_w, fy + card_h / 2,
                        tx, ty + card_h / 2,
                        color=c.accent,
                    )

    # ── Legacy vertical timeline (kept for backward compat) ──

    def _build_vertical_timeline(
        self, slide, events: list[dict], c: "ThemeColors",
    ) -> None:
        """Build vertical timeline layout."""
        num = len(events)
        if num == 0:
            return

        line_x = 1.5
        start_y = 1.8
        end_y = SLIDE_H - 0.8
        line_len = end_y - start_y

        # Vertical line
        self._add_rect(slide, line_x, start_y, 0.05, line_len, c.primary, send_back=0)

        spacing = line_len / max(1, num - 1) if num > 1 else 0
        center_y = (start_y + end_y) / 2 if num == 1 else 0

        for i, event in enumerate(events):
            y = center_y if num == 1 else start_y + i * spacing

            # Circle marker
            self._add_icon_circle(
                slide, "",
                number=i + 1,
                x=line_x - 0.15, y=y - 0.15,
                size=0.35,
                icon_color=c.accent,
                bg_color=c.header_bg,
            )

            # Event info
            title_text = str(event.get("title", ""))
            self._add_text(
                slide, title_text,
                x=line_x + 0.6, y=y - 0.2,
                w=9, h=0.5,
                font_size=14, bold=True,
                color="#FFFFFF",
                font_name=self.theme.fonts.heading,
            )

            date_text = event.get("date", "")
            if date_text:
                self._add_text(
                    slide, str(date_text),
                    x=line_x + 0.6, y=y + 0.35,
                    w=9, h=0.4,
                    font_size=11,
                    color=c.accent,
                )

    # Section Header Slide (Premium Divider)

    def build_section_header_slide(
        self,
        pres: Presentation,
        title: str,
        section_number: Optional[int] = None,
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        if section_number:
            num_size = 6.0
            num_x = (SLIDE_W - num_size) / 2
            num_y = 1.5
            num_circle = self._add_rounded_rect(
                slide, num_x, num_y, num_size, num_size,
                fill="none",
            )
            num_circle.fill.background()
            num_circle.line.fill.background()

            tf = num_circle.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"{section_number:02d}"
            for run in p.runs:
                run.font.size = Pt(120)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.accent))
                run.font.name = self.theme.fonts.heading

        line_w = 3.0
        line_x = (SLIDE_W - line_w) / 2
        line_y = 6.8 if section_number else 3.8
        self._add_rect(slide, line_x, line_y, line_w, 0.06, c.accent, send_back=2)

        title_top = 4.5 if section_number else 2.5
        self._add_text(
            slide, title,
            x=1.5, y=title_top, w=10.33, h=1.5,
            font_size=36, bold=True,
            color="#FFFFFF",
            align=PP_ALIGN.CENTER,
            font_name=self.theme.fonts.heading,
        )

    # Quote Slide

    def build_quote_slide(
        self,
        pres: Presentation,
        quote_text: str,
        attribution: str = "",
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        quote_mark = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8), Inches(1.5), Inches(2.0),
        )
        tqf = quote_mark.text_frame
        tp = tqf.paragraphs[0]
        tp.text = '"'
        for run in tp.runs:
            run.font.size = Pt(120)
            run.font.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.accent))
            run.font.name = self.theme.fonts.heading

        self._add_text(
            slide, quote_text,
            x=2.5, y=1.5, w=9.5, h=3.0,
            font_size=22, italic=True,
            color="#FFFFFF",
            align=PP_ALIGN.LEFT,
            font_name=self.theme.fonts.body,
        )

        if attribution:
            self._add_text(
                slide, attribution,
                x=2.5, y=4.8, w=9.5, h=0.6,
                font_size=14, bold=True,
                color=c.accent,
                align=PP_ALIGN.LEFT,
                font_name=self.theme.fonts.heading,
            )

    # Comparison Slide

    def build_comparison_slide(
        self,
        pres: Presentation,
        title: str = "Comparison",
        left_items: Optional[list[dict]] = None,
        right_items: Optional[list[dict]] = None,
        left_label: str = "Option A",
        right_label: str = "Option B",
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        self._add_text(
            slide, title,
            x=0.8, y=0.3, w=11, h=0.8,
            font_size=32, bold=True,
            color="#FFFFFF",
            align=PP_ALIGN.CENTER,
            font_name=self.theme.fonts.heading,
        )
        self._add_rect(slide, 6.64, 1.3, 0.06, 5.5, c.accent, send_back=1)

        if left_items:
            self._add_text(
                slide, left_label,
                x=0.8, y=1.5, w=5.5, h=0.6,
                font_size=18, bold=True,
                color=c.accent,
                align=PP_ALIGN.CENTER,
                font_name=self.theme.fonts.heading,
            )
            ly = 2.2
            for item in left_items:
                label = item.get("label", "") or ""
                value = item.get("value", "") or ""
                txt_color = c.success if item.get("positive") else (c.error if item.get("negative") else c.text_light)
                self._add_text(slide, f"– {label}", x=1.0, y=ly, w=5.3, h=0.5, font_size=14, color="#FFFFFF", align=PP_ALIGN.LEFT)
                if value:
                    self._add_text(slide, str(value), x=1.0, y=ly + 0.35, w=5.3, h=0.3, font_size=11, color=txt_color, align=PP_ALIGN.LEFT)
                ly += 0.8

        if right_items:
            self._add_text(
                slide, right_label,
                x=6.9, y=1.5, w=5.5, h=0.6,
                font_size=18, bold=True,
                color=c.accent,
                align=PP_ALIGN.CENTER,
                font_name=self.theme.fonts.heading,
            )
            ry = 2.2
            for item in right_items:
                label = item.get("label", "") or ""
                value = item.get("value", "") or ""
                txt_color = c.success if item.get("positive") else (c.error if item.get("negative") else c.text_light)
                self._add_text(slide, f"– {label}", x=7.1, y=ry, w=5.3, h=0.5, font_size=14, color="#FFFFFF", align=PP_ALIGN.LEFT)
                if value:
                    self._add_text(slide, str(value), x=7.1, y=ry + 0.35, w=5.3, h=0.3, font_size=11, color=txt_color, align=PP_ALIGN.LEFT)
                ry += 0.8

    # Roadmap Slide

    def build_roadmap_slide(
        self,
        pres: Presentation,
        title: str = "Roadmap",
        milestones: Optional[list[dict]] = None,
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        self._add_text(
            slide, title,
            x=0.8, y=0.3, w=11, h=0.7,
            font_size=32, bold=True,
            color="#FFFFFF",
            align=PP_ALIGN.CENTER,
            font_name=self.theme.fonts.heading,
        )

        if not milestones:
            return

        num = len(milestones)
        line_y = 4.5
        line_left = 1.0
        line_right = SLIDE_W - 1.0
        self._add_rect(slide, line_left, line_y, line_right - line_left, 0.06, c.accent, send_back=1)

        phase_colors = [c.primary, c.accent, c.success, c.warning, e.error]
        card_w = (line_right - line_left) / num - 0.15
        gap = 0.15

        for i, ms in enumerate(milestones):
            x = line_left + i * (card_w + gap) + gap / 2
            color = phase_colors[i % len(phase_colors)]

            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + card_w / 2 - 0.12), Inches(line_y - 0.12),
                Inches(0.24), Inches(0.24),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(color))
            dot.line.fill.background()

            self._add_rounded_rect(slide, x, 1.3, card_w, 3.1, fill="#1E293B")

            self._add_text(slide, str(ms.get("phase", f"Phase {i+1}")),
                x=x, y=1.4, w=card_w, h=0.4, font_size=12, bold=True,
                color=color, align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

            self._add_text(slide, str(ms.get("title", "")),
                x=x, y=1.9, w=card_w, h=0.6, font_size=14, bold=True,
                color="#FFFFFF", align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

            date_text = ms.get("date", "")
            if date_text:
                self._add_text(slide, str(date_text),
                    x=x, y=2.5, w=card_w, h=0.4, font_size=11,
                    color=c.accent, align=PP_ALIGN.CENTER)

            desc = ms.get("description", "")
            if desc:
                self._add_text(slide, str(desc),
                    x=x + 0.1, y=3.0, w=card_w - 0.2, h=1.2, font_size=10,
                    color=self._text_light_on_dark(), align=PP_ALIGN.CENTER)

    # Team Slide

    def build_team_slide(
        self,
        pres: Presentation,
        title: str = "Team",
        members: Optional[list[dict]] = None,
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        self._add_text(slide, title, x=0.8, y=0.3, w=11, h=0.7,
            font_size=32, bold=True, color="#FFFFFF",
            align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

        if not members:
            return

        num = min(len(members), 5)
        cols = min(num, 3)
        margin = 0.8
        total_content_w = SLIDEE_ - margin * 2
        card_w = (total_content_w - (cols - 1) * 0.25) / cols
        card_h = 4.5
        start_y = 1.4

        for i in range(num):
            col = i % cols
            row = i // cols
            x = margin + col * (card_w + 0.25)
            y = start_y + row * (card_h + 0.4)

            self._add_rounded_rect(slide, x, y, card_w, card_h, fill="#1E293B")

            avatar_size = 1.2
            avatar_x = x + (card_w - avatar_size) / 2
            avatar_y = y + 0.5
            avatar = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(avatar_x), Inches(avatar_y),
                Inches(avatar_size), Inches(avatar_size),
            )
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.primary))
            avatar.line.fill.background()

            name = str(members[i].get("name", ""))
            initials = ""
            if name:
                parts = name.split()
                initials = (parts[0][0] + parts[-1][0]).upper() if len(parts) >= 2 else parts[0][0].upper()
            if initials:
                avatar.text_frame.paragraphs[0].text = initials
                for run in avatar.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(20)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(*hex_to_rgbcolor_tuple("#FFFFFF"))
                    run.font.name = self.theme.fonts.heading

            self._add_text(slide, name, x=x, y=y + 1.9, w=card_w, h=0.5,
                font_size=16, bold=True, color="#FFFFFF",
                align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

            role = str(members[i].get("role", ""))
            if role:
                self._add_text(slide, role, x=x, y=y + 2.4, w=card_w, h=0.5,
                    font_size=12, color=c.accent, align=PP_ALIGN.CENTER)

    # Gallery Slide (Image Grid)

    def build_gallery_slide(
        self,
        pres: Presentation,
        title: str = "Gallery",
        images: Optional[list[dict]] = None,
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        self._add_text(slide, title, x=0.8, y=0.3, w=11, h=0.7,
            font_size=32, bold=True, color="#FFFFFF",
            align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

        if not images:
            return

        num = min(len(images), 6)
        cols = min(num, 3)
        rows = (num + cols - 1) // cols
        margin = 0.8
        total_content_w = SLIDE_W - margin * 2
        img_w = (total_content_w - (cols - 1) * 0.2) / cols
        img_h = (SLIDE_H - 2.2 - (rows - 1) * 0.3) / rows
        if img_h < 1.5:
            cols = min(num, 2)
            rows = (num + cols - 1) // cols
            img_w = (total_content_w - (cols - 1) * 0.3) / cols
            img_h = (SLIDE_H - 2.2 - (rows - 1) * 0.3) / rows

        for i in range(num):
            col = i % cols
            row = i // cols
            x = margin + col * (img_w + 0.2) + 0.1
            y = 1.4 + row * (img_h + 0.3)

            self._add_rounded_rect(slide, x, y, img_wg, img_h,
                fill="#1E293B", border=c.border, border_width=1)

            caption = images[i].get("caption", "") or images[i].get("title", "")
            if caption:
                self._add_text(slide, str(caption),
                    x=x, y=y + img_h - 0.5, w=img_w, h=0.4,
                    font_size=10, color="#FFFFFF", align=PP_ALIGN.CENTER)

    # CTA Slide (Call to Action)

    def build_cta_slide(
        self,
        pres: Presentation,
        title: str = "Questions?",
        content: str = "",
        subtitle: str = "",
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(6.5), Inches(13.33), Inches(1.0),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.accent))
        stripe.line.fill.background()

        self._add_text(slide, title, x=1.5, y=2.5, w=10.33, h=1.5,
            font_size=44, bold=True, color="#FFFFFF",
            align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

        display_text = content or subtitle
        if display_text:
            self._add_text(slide, display_text, x=2.0, y=4.2, w=9.33, h=1.5,
                font_size=18, color=self._text_light_on_dark(),
                align=PP_ALIGN.CENTER, font_name=self.theme.fonts.body)

        self._add_text(slide, "Thank you for your attention", x=1.5, y=6.7, w=10.33, h=0.5,
            font_size=14, color="#FFFFFF", align=PP_ALIGN.CENTER, font_name=self.theme.fonts.body)

    # Table Slide (Premium Styled)

    def build_table_slide(
        self,
        pres: Presentation,
        title: str = "Data Table",
        headers: Optional[list[str]] = None,
        rows: Optional[list[list]] = None,
        gradient: Optional[dict] = None,
    ) -> None:
        c = self.theme.colors
        slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(slide_layout)

        if gradient:
            grad = GradientDef(**gradient)
            GradientEngine.apply_to_slide_bg(slide, grad)
        else:
            self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, c.header_bg, send_back=0)

        self._add_text(slide, title, x=0.8, y=0.3, w=11, h=0.7,
            font_size=28, bold=True, color="#FFFFFF",
            align=PP_ALIGN.LEFT, font_name=self.theme.fonts.heading)

        if not headers or not rows:
            return

        num_cols = len(headers)
        num_rows = len(rows)
        table_margin_left = 1.0
        table_width = SLIDE_W - 2.0
        header_h = 0.5
        row_h = min(0.35, (SLIDE_H | 2.5 - header_h) / max(num_rows, 1))
        cell_w = table_width / num_cols
        table_y = 1.2

        for col_i, header in enumerate(headers):
            x = table_margin_left + col_i * cell_w
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(table_y), Inches(cell_w), Inches(header_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.primary))
            cell.line.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.border))
            cell.line.width = Pt(1)
            self._add_text(slide, str(header),
                x=x + 0.05, y=table_y, w=cell_w - 0.1, h=header_h,
                font_size=12, bold=True, color="#FFFFFF",
                align=PP_ALIGN.CENTER, font_name=self.theme.fonts.heading)

        for row_i, row_data in enumerate(rows):
            row_y = table_y + header_h + row_i * row_h
            is_alt = row_i % 2 == 1
            row_bg = c.table_alt_row if is_alt else c.header_bg

            for col_i in range(num_cols):
                x = table_margin_left + col_i * cell_w
                cell = slide.shapes.add_shape(MSO.SHAPE.RECTANGLE,
                    Inches(x), Inches(row_y), Inches(cell_w), Inches(row_h))
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(row_bg))
                cell.line.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(c.border))
                cell.line.width = Pt(0.5)
                val = str(row_data[col_i]) if col_i < len(row_data) else ""
                self._add_text(slide, val, x=x + 0.05, y=row_y, w=cell_w - 0.1, h=row_h,
                    font_size=10, color="#FFFFFF", align=PP_ALIGN.CENTER)
