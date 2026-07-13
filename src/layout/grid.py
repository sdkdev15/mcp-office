"""Grid system for auto-layout of items in presentations and documents.

Supports:
- Fixed grid layouts (2x2, 3x3, 4x2, etc)
- Masonry layouts for varying height items
- Responsive grid with automatic column calculation
- Card-based layouts with spacing and alignment
"""

from __future__ import annotations

from typing import Optional
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from src.styles.themes import get_theme
from src.utils.colors import hex_to_rgbcolor_tuple
from src.utils.logger import get_logger

log = get_logger("grid")


class GridItem:
    """Represents a single item in a grid layout."""
    
    def __init__(self, title: str = "", subtitle: str = "", content: str = "", 
                 icon: Optional[str] = None, image: Optional[str] = None,
                 color: Optional[str] = None, width: float = None, height: float = None):
        self.title = title
        self.subtitle = subtitle
        self.content = content
        self.icon = icon
        self.image = image
        self.color = color
        self.width = width
        self.height = height


class GridLayout:
    """Create grid layouts for presentations."""
    
    @staticmethod
    def fixed_grid(slide, items: list[dict], grid_config: dict, left: float = 1.0, 
                   top: float = 1.5, theme_name: str = "corporate") -> None:
        """Create a fixed grid layout.
        
        Args:
            slide: PowerPoint slide object
            items: List of item dicts with title, subtitle, content, icon, etc.
            grid_config: Dict with columns, rows, spacing, etc.
            left: Left margin in inches
            top: Top margin in inches
            theme_name: Theme name for colors
        """
        theme = get_theme(theme_name)
        colors = theme.colors
        
        cols = grid_config.get("columns", 3)
        rows = grid_config.get("rows", 2)
        spacing = grid_config.get("spacing", 0.2)
        card_width = grid_config.get("card_width")
        card_height = grid_config.get("card_height")
        
        if not card_width or not card_height:
            # Calculate based on available space
            total_width = slide.slide_width / 914400 - left * 2  # Convert EMU to inches
            total_height = slide.slide_height / 914400 - top * 2
            
            card_width = (total_width - (cols - 1) * spacing) / cols
            card_height = (total_height - (rows - 1) * spacing) / rows
        
        for i, item_data in enumerate(items):
            if i >= cols * rows:
                break
            
            col = i % cols
            row = i // cols
            
            x = left + col * (card_width + spacing)
            y = top + row * (card_height + spacing)
            
            # Draw card background
            card = slide.shapes.add_shape(
                __import__('pptx.enum.shapes', fromlist=['MSO_SHAPE']).MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(card_width), Inches(card_height),
            )
            card.fill.solid()
            
            # Use item color if provided, otherwise use theme colors
            if item_data.get("color"):
                color = hex_to_rgbcolor_tuple(item_data["color"])
            else:
                color = hex_to_rgbcolor_tuple(colors.primary)
            
            card.fill.fore_color.rgb = RGBColor(*color)
            card.line.fill.background()
            
            # Add title
            title_text = item_data.get("title", "")
            if title_text:
                slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(y + 0.1),
                    Inches(card_width - 0.2), Inches(0.5),
                ).text_frame.paragraphs[0].text = title_text
            
            # Add subtitle
            subtitle_text = item_data.get("subtitle", "")
            if subtitle_text:
                slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(y + 0.5),
                    Inches(card_width - 0.2), Inches(0.3),
                ).text_frame.paragraphs[0].text = subtitle_text
            
            # Add content
            content_text = item_data.get("content", "")
            if content_text:
                slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(y + 0.8),
                    Inches(card_width - 0.2), Inches(card_height - 1.0),
                ).text_frame.paragraphs[0].text = content_text
    
    @staticmethod
    def masonry_grid(slide, items: list[dict], columns: int = 3, left: float = 1.0,
                    top: float = 1.5, theme_name: str = "corporate") -> None:
        """Create a masonry-style grid layout.
        
        Args:
            slide: PowerPoint slide object
            items: List of item dicts
            columns: Number of columns
            left: Left margin in inches
            top: Top margin in inches
            theme_name: Theme name for colors
        """
        theme = get_theme(theme_name)
        colors = theme.colors
        
        # Calculate column width
        total_width = slide.slide_width / 914400 - left * 2
        card_width = (total_width - (columns - 1) * 0.2) / columns
        spacing = 0.2
        
        # Track row heights for each column
        col_heights = [top] * columns
        
        for i, item_data in enumerate(items):
            col = i % columns
            
            # Calculate y position
            y = col_heights[col]
            
            # Determine card height based on content length
            content = item_data.get("content", "")
            title = item_data.get("title", "")
            height = 1.5 if len(content) < 50 else (2.0 if len(content) < 100 else 2.5)
            
            # Draw card
            card = slide.shapes.add_shape(
                __import__('pptx.enum.shapes', fromlist=['MSO_SHAPE']).MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + col * (card_width + spacing)), Inches(y),
                Inches(card_width), Inches(height),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(colors.surface))
            card.line.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(colors.border))
            card.line.width = Pt(1)
            
            # Add content to card
            if title:
                slide.shapes.add_textbox(
                    Inches(left + col * (card_width + spacing) + 0.1),
                    Inches(y + 0.1),
                    Inches(card_width - 0.2), Inches(0.5),
                ).text_frame.paragraphs[0].text = title
            
            if content:
                slide.shapes.add_textbox(
                    Inches(left + col * (card_width + spacing) + 0.1),
                    Inches(y + 0.6),
                    Inches(card_width - 0.2), Inches(height - 0.7),
                ).text_frame.paragraphs[0].text = content
            
            # Update column height
            col_heights[col] += height + spacing
    
    @staticmethod
    def card_layout(slide, items: list[dict], columns: int = 3, left: float = 1.0,
                   top: float = 1.5, card_width: float = 3.0, card_height: float = 1.5,
                   theme_name: str = "corporate") -> None:
        """Create a simple card layout.
        
        Args:
            slide: PowerPoint slide object
            items: List of item dicts
            columns: Number of columns
            left: Left margin in inches
            top: Top margin in inches
            card_width: Width of each card in inches
            card_height: Height of each card in inches
            theme_name: Theme name for colors
        """
        theme = get_theme(theme_name)
        colors = theme.colors
        spacing = 0.2
        
        for i, item_data in enumerate(items):
            col = i % columns
            
            x = left + col * (card_width + spacing)
            y = top + (i // columns) * (card_height + spacing)
            
            # Draw card
            card = slide.shapes.add_shape(
                __import__('pptx.enum.shapes', fromlist=['MSO_SHAPE']).MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(card_width), Inches(card_height),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*hex_to_rgbcolor_tuple(colors.surface))
            card.line.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(colors.border))
            card.line.width = Pt(1)
            
            # Add title
            if item_data.get("title"):
                slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(y + 0.1),
                    Inches(card_width - 0.2), Inches(0.5),
                ).text_frame.paragraphs[0].text = item_data["title"]
            
            # Add subtitle
            if item_data.get("subtitle"):
                slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(y + 0.6),
                    Inches(card_width - 0.2), Inches(0.3),
                ).text_frame.paragraphs[0].text = item_data["subtitle"]
            
            # Add content
            if item_data.get("content"):
                slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(y + 0.9),
                    Inches(card_width - 0.2), Inches(card_height - 1.0),
                ).text_frame.paragraphs[0].text = item_data["content"]
