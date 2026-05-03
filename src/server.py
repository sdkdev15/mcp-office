"""MCP Office Server — Main entry point for mcp 1.27.x."""

from __future__ import annotations

import asyncio
import os
import uuid
from typing import Any, Optional

from mcp.server import Server
from mcp.server.models import InitializationOptions
from mcp.types import TextContent, Tool, ServerCapabilities
from mcp.server.stdio import stdio_server
from mcp.server.sse import SseServerTransport

from src.generators.excel_generator import ExcelGenerator
from src.generators.docx_generator import DOCXGenerator
from src.generators.pptx_generator import PPTXGenerator
from src.generators.odf_generator import ODFGenerator
from src.styles.themes import list_themes
from src.utils.file_handler import FileHandler
from src.utils.cleanup import FileCleanup
from src.utils.rate_limiter import RateLimiter
from src.utils.logger import get_logger
from src.utils.validators import ValidationError

# ── Initialize Components ──

log = get_logger("server")
app = Server("mcp-office", version="1.0.0")
file_handler = FileHandler(output_dir=os.environ.get("OUTPUT_DIR", "outputs"))
file_cleanup = FileCleanup(output_dir=os.environ.get("OUTPUT_DIR", "outputs"))
rate_limiter = RateLimiter()
base_url = os.environ.get("BASE_URL", "")  # e.g., "http://mcp-office:8765" for SSE mode


def get_session_id(params: dict[str, Any]) -> str:
    return params.get("session_id") or str(uuid.uuid4())[:8]


def check_rate_limit(session_id: str) -> Optional[str]:
    allowed, info = rate_limiter.is_allowed(session_id)
    if not allowed:
        return f"Rate limit exceeded. Try again in {info.get('retry_after', 60)}s."
    return None


def format_file_result(file_info: dict, base_url: str = "") -> str:
    result = (
        f"✅ **File Generated Successfully**\n\n"
        f"**File:** {file_info['filename']}\n"
        f"**Size:** {file_handler._human_readable_size(file_info['size_bytes'])}\n"
        f"**Type:** {file_info['mime_type']}\n"
        f"**Path:** {file_info['filepath']}\n"
        f"**Created:** {file_info['created_at']}"
    )
    if base_url:
        download_url = f"{base_url}/files/{file_info['session_id']}/{file_info['filename']}"
        result += f"\n\n**Download:** [{download_url}]({download_url})"
    return result


# ── Tool Definitions ──

TOOLS = [
    Tool(
        name="excel_create",
        description="Create an Excel workbook (.xlsx) with multiple sheets, styling, and formatting.",
        inputSchema={
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "Output filename (e.g., 'report.xlsx')"},
                "sheets": {
                    "type": "array",
                    "description": "List of sheet objects. Each sheet should have 'name' (string), 'headers' (array of strings), and 'rows' (array of arrays).",
                    "items": {"type": "object"},
                },
                "theme": {"type": "string", "description": "Theme name (corporate, minimal, creative, academic, dark)", "default": "corporate"},
                "template_path": {"type": "string", "description": "Optional path to a .xlsx template file to use as base (preserves formatting, formulas, charts)"},
                "session_id": {"type": "string", "description": "Optional session ID"},
                "metadata": {"type": "object", "description": "Optional document metadata"},
                "locale": {"type": "string", "description": "Locale (en_US, id_ID)", "default": "en_US"},
            },
            "required": ["filename", "sheets"],
        },
    ),
    Tool(
        name="excel_export",
        description="Export data to multiple formats (xlsx, ods, csv).",
        inputSchema={
            "type": "object",
            "properties": {
                "sheets": {
                    "type": "array",
                    "description": "List of sheet objects with name, headers, and rows",
                    "items": {"type": "object"},
                },
                "format": {"type": "string", "description": "xlsx, ods, csv, all", "default": "all"},
                "session_id": {"type": "string"},
            },
            "required": ["sheets"],
        },
    ),
    Tool(
        name="docx_create",
        description="Create a Word document (.docx) with page setup and theme styling.",
        inputSchema={
            "type": "object",
            "properties": {
                "filename": {"type": "string"},
                "title": {"type": "string", "default": "Document"},
                "theme": {"type": "string", "default": "corporate"},
                "template_path": {"type": "string", "description": "Optional path to a .docx template file to use as base (preserves formatting, styles, headers, footers)"},
                "page_size": {"type": "string", "default": "A4"},
                "orientation": {"type": "string", "default": "portrait"},
                "content_paragraphs": {"type": "array", "description": "Optional list of paragraph texts to append", "items": {"type": "string"}},
                "tables": {"type": "array", "description": "Optional list of table data with headers and rows", "items": {"type": "object"}},
                "session_id": {"type": "string"},
                "metadata": {"type": "object"},
            },
            "required": ["filename"],
        },
    ),
    Tool(
        name="docx_generate_from_prompt",
        description="Generate a Word document from a natural language description.",
        inputSchema={
            "type": "object",
            "properties": {
                "prompt": {"type": "string"},
                "filename": {"type": "string"},
                "theme": {"type": "string", "default": "corporate"},
                "session_id": {"type": "string"},
            },
            "required": ["prompt"],
        },
    ),
    Tool(
        name="docx_export",
        description="Export a document in multiple formats (docx, odt).",
        inputSchema={
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "format": {"type": "string", "default": "all"},
                "theme": {"type": "string", "default": "corporate"},
                "session_id": {"type": "string"},
                "metadata": {"type": "object"},
            },
            "required": ["title"],
        },
    ),
    Tool(
        name="pptx_create",
        description="Create a PowerPoint presentation (.pptx) with slides and theme styling.",
        inputSchema={
            "type": "object",
            "properties": {
                "filename": {"type": "string"},
                "title": {"type": "string", "default": "Presentation"},
                "theme": {"type": "string", "default": "corporate"},
                "template_path": {"type": "string", "description": "Optional path to a .pptx template file to use as base (preserves master slides, themes, layouts)"},
                "slides": {
                    "type": "array",
                    "description": "Optional list of slide data with title, content, bullets",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide title"},
                            "content": {"type": "string", "description": "Slide body content/description"},
                            "bullets": {"type": "array", "items": {"type": "string"}, "description": "Bullet points for the slide"},
                            "image_path": {"type": "string", "description": "Optional path to an image"},
                            "table_headers": {"type": "array", "items": {"type": "string"}, "description": "Optional table headers"},
                            "table_rows": {"type": "array", "items": {"type": "array"}, "description": "Optional table rows"}
                        }
                    }
                },
                "slide_size": {"type": "string", "default": "widescreen"},
                "session_id": {"type": "string"},
                "metadata": {"type": "object"},
            },
            "required": ["filename"],
        },
    ),
    Tool(
        name="pptx_generate_from_prompt",
        description="Generate a PowerPoint presentation from a natural language description.",
        inputSchema={
            "type": "object",
            "properties": {
                "prompt": {"type": "string"},
                "filename": {"type": "string"},
                "theme": {"type": "string", "default": "corporate"},
                "session_id": {"type": "string"},
            },
            "required": ["prompt"],
        },
    ),
    Tool(
        name="pptx_export",
        description="Export a presentation in multiple formats (pptx, odp).",
        inputSchema={
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "format": {"type": "string", "default": "all"},
                "theme": {"type": "string", "default": "corporate"},
                "session_id": {"type": "string"},
                "metadata": {"type": "object"},
            },
            "required": ["title"],
        },
    ),
    Tool(
        name="list_themes_tool",
        description="List all available themes for document generation.",
        inputSchema={"type": "object", "properties": {}},
    ),
    Tool(
        name="list_files",
        description="List all generated files for a session.",
        inputSchema={
            "type": "object",
            "properties": {"session_id": {"type": "string"}},
            "required": ["session_id"],
        },
    ),
    Tool(
        name="get_storage_stats",
        description="Get current storage usage statistics.",
        inputSchema={"type": "object", "properties": {}},
    ),
]


# ── MCP Handlers ──

@app.list_tools()
async def list_tools() -> list[Tool]:
    """List available tools."""
    return TOOLS


@app.call_tool()
async def call_tool(name: str, arguments: dict[str, Any]) -> list[TextContent]:
    """Handle tool calls."""
    session_id = get_session_id(arguments)

    if name not in ("list_themes_tool", "list_files", "get_storage_stats"):
        rate_error = check_rate_limit(session_id)
        if rate_error:
            return [TextContent(type="text", text=rate_error)]

    try:
        if name == "excel_create":
            return await _excel_create(arguments)
        elif name == "excel_export":
            return await _excel_export(arguments)
        elif name == "docx_create":
            return await _docx_create(arguments)
        elif name == "docx_generate_from_prompt":
            return await _docx_generate_from_prompt(arguments)
        elif name == "docx_export":
            return await _docx_export(arguments)
        elif name == "pptx_create":
            return await _pptx_create(arguments)
        elif name == "pptx_generate_from_prompt":
            return await _pptx_generate_from_prompt(arguments)
        elif name == "pptx_export":
            return await _pptx_export(arguments)
        elif name == "list_themes_tool":
            return await _list_themes()
        elif name == "list_files":
            return await _list_files(arguments)
        elif name == "get_storage_stats":
            return await _get_storage_stats()
        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]

    except ValidationError as e:
        return [TextContent(type="text", text=f"Validation Error: {e.message}")]
    except Exception as e:
        log.error(f"Tool {name} failed: {e}")
        return [TextContent(type="text", text=f"Error: {str(e)}")]


# ── Tool Implementations ──

async def _excel_create(args: dict) -> list[TextContent]:
    gen = ExcelGenerator(args.get("theme", "corporate"))
    template_path = args.get("template_path")
    
    if template_path:
        data = gen.create_from_template(
            template_path,
            args["sheets"],
            args.get("metadata"),
        )
    else:
        data = gen.create_workbook(args["sheets"], args.get("metadata"))
    
    filename = args["filename"]
    if not filename.endswith(".xlsx"):
        filename += ".xlsx"
    filename = file_handler.generate_filename(filename.rsplit(".", 1)[0], ".xlsx")
    file_info = await file_handler.save_file(data, filename, get_session_id(args))
    return [TextContent(type="text", text=format_file_result(file_info, base_url))]


async def _excel_export(args: dict) -> list[TextContent]:
    results = []
    fmt = args.get("format", "all").lower()
    session_id = get_session_id(args)

    if fmt in ("xlsx", "all"):
        gen = ExcelGenerator()
        data = gen.create_workbook(args["sheets"])
        filename = file_handler.generate_filename("export", ".xlsx")
        file_info = await file_handler.save_file(data, filename, session_id)
        results.append(format_file_result(file_info))

    if fmt in ("ods", "all"):
        gen = ODFGenerator()
        data = gen.create_spreadsheet(args["sheets"])
        filename = file_handler.generate_filename("export", ".ods")
        file_info = await file_handler.save_file(data, filename, session_id)
        results.append(format_file_result(file_info))

    return [TextContent(type="text", text="\n\n".join(results))]


async def _docx_create(args: dict) -> list[TextContent]:
    gen = DOCXGenerator(args.get("theme", "corporate"))
    template_path = args.get("template_path")
    content_paragraphs = args.get("content_paragraphs")
    tables = args.get("tables")
    
    if template_path:
        data = gen.create_from_template(
            template_path,
            title=args.get("title", "Document"),
            content_paragraphs=content_paragraphs,
            tables=tables,
            metadata=args.get("metadata"),
        )
    else:
        from docx import Document
        gen.doc = Document()
        gen._setup_page(args.get("page_size", "A4"), args.get("orientation", "portrait"))
        gen._apply_theme_styles()
        
        if args.get("metadata"):
            gen._apply_metadata(args["metadata"])
        
        title = args.get("title", "Document")
        gen.add_heading(title, level=1)
        gen.add_horizontal_line()
        
        if content_paragraphs:
            combined = "\n".join(content_paragraphs)
            gen._parse_and_render(combined)
        
        if tables:
            for table_data in tables:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                gen.add_table(headers, rows)
        
        data = gen._save_document()
    
    filename = args["filename"]
    if not filename.endswith(".docx"):
        filename += ".docx"
    filename = file_handler.generate_filename(filename.rsplit(".", 1)[0], ".docx")
    file_info = await file_handler.save_file(data, filename, get_session_id(args))
    return [TextContent(type="text", text=format_file_result(file_info, base_url))]


async def _docx_generate_from_prompt(args: dict) -> list[TextContent]:
    gen = DOCXGenerator(args.get("theme", "corporate"))
    title = args.get("filename") or "Generated Document"
    data = gen.create_from_prompt(args["prompt"], title)
    filename = args.get("filename") or "document.docx"
    if not filename.endswith(".docx"):
        filename += ".docx"
    filename = file_handler.generate_filename(filename.rsplit(".", 1)[0], ".docx")
    file_info = await file_handler.save_file(data, filename, get_session_id(args))
    return [TextContent(type="text", text=format_file_result(file_info, base_url))]


async def _docx_export(args: dict) -> list[TextContent]:
    results = []
    fmt = args.get("format", "all").lower()
    session_id = get_session_id(args)
    title = args.get("title", "Document")

    if fmt in ("docx", "all"):
        gen = DOCXGenerator(args.get("theme", "corporate"))
        data = gen.create_document(title, metadata=args.get("metadata"))
        filename = file_handler.generate_filename(title, ".docx")
        file_info = await file_handler.save_file(data, filename, session_id)
        results.append(format_file_result(file_info))

    if fmt in ("odt", "all"):
        gen = ODFGenerator()
        odt_doc = gen.create_text_document(title, args.get("metadata"))
        gen.add_paragraph_to_odt(odt_doc, f"Generated: {title}")
        data = gen.save_odt(odt_doc)
        filename = file_handler.generate_filename(title, ".odt")
        file_info = await file_handler.save_file(data, filename, session_id)
        results.append(format_file_result(file_info))

    return [TextContent(type="text", text="\n\n".join(results))]


async def _pptx_create(args: dict) -> list[TextContent]:
    gen = PPTXGenerator(args.get("theme", "corporate"))
    template_path = args.get("template_path")
    slides = args.get("slides")
    
    if template_path:
        data = gen.create_from_template(
            template_path,
            slides=slides,
            metadata=args.get("metadata"),
        )
    else:
        # Create presentation from scratch
        from pptx import Presentation
        from pptx.util import Inches
        gen.pres = Presentation()
        
        # Set slide size
        slide_size = args.get("slide_size", "widescreen")
        if slide_size.lower() == "standard":
            gen.pres.slide_width = Inches(10)
            gen.pres.slide_height = Inches(7.5)
        else:
            gen.pres.slide_width = Inches(13.33)
            gen.pres.slide_height = Inches(7.5)
        
        # Apply metadata
        if args.get("metadata"):
            gen._apply_metadata(args["metadata"])
        
        # Add title slide
        title = args.get("title", "Presentation")
        gen.add_slide("title", title=title)
        
        # Add content slides
        if slides:
            for slide_data in slides:
                gen.add_slide(
                    layout=slide_data.get("layout", "title_and_content"),
                    title=slide_data.get("title"),
                    content=slide_data.get("content"),
                    bullets=slide_data.get("bullets"),
                )
        
        data = gen._save_presentation()
    
    filename = args["filename"]
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    filename = file_handler.generate_filename(filename.rsplit(".", 1)[0], ".pptx")
    file_info = await file_handler.save_file(data, filename, get_session_id(args))
    return [TextContent(type="text", text=format_file_result(file_info, base_url))]


async def _pptx_generate_from_prompt(args: dict) -> list[TextContent]:
    gen = PPTXGenerator(args.get("theme", "corporate"))
    title = args.get("filename") or "Generated Presentation"
    data = gen.create_from_prompt(args["prompt"], title)
    filename = args.get("filename") or "presentation.pptx"
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    filename = file_handler.generate_filename(filename.rsplit(".", 1)[0], ".pptx")
    file_info = await file_handler.save_file(data, filename, get_session_id(args))
    return [TextContent(type="text", text=format_file_result(file_info, base_url))]


async def _pptx_export(args: dict) -> list[TextContent]:
    results = []
    fmt = args.get("format", "all").lower()
    session_id = get_session_id(args)
    title = args.get("title", "Presentation")

    if fmt in ("pptx", "all"):
        gen = PPTXGenerator(args.get("theme", "corporate"))
        data = gen.create_presentation(title, metadata=args.get("metadata"))
        filename = file_handler.generate_filename(title, ".pptx")
        file_info = await file_handler.save_file(data, filename, session_id)
        results.append(format_file_result(file_info))

    if fmt in ("odp", "all"):
        gen = ODFGenerator()
        odp_doc = gen.create_presentation(title, args.get("metadata"))
        data = gen.save_odp(odp_doc)
        filename = file_handler.generate_filename(title, ".odp")
        file_info = await file_handler.save_file(data, filename, session_id)
        results.append(format_file_result(file_info))

    return [TextContent(type="text", text="\n\n".join(results))]


async def _list_themes() -> list[TextContent]:
    themes = list_themes()
    result = "**Available Themes:**\n\n"
    for theme in themes:
        result += f"- **{theme['name']}**: {theme['description']}\n  - Primary: {theme['primary_color']}\n\n"
    return [TextContent(type="text", text=result)]


async def _list_files(args: dict) -> list[TextContent]:
    session_id = args.get("session_id", "")
    files = file_handler.list_session_files(session_id)
    if not files:
        return [TextContent(type="text", text=f"No files found for session: {session_id}")]
    result = f"**Files in session {session_id}:**\n\n"
    for f in files:
        result += f"- `{f['filename']}` ({f['size_bytes']} bytes)\n  - Created: {f['created_at']}\n\n"
    return [TextContent(type="text", text=result)]


async def _get_storage_stats() -> list[TextContent]:
    stats = file_cleanup.get_storage_stats()
    result = (
        f"**Storage Statistics:**\n\n"
        f"- Total files: {stats['total_files']}\n"
        f"- Total size: {stats['total_size_human']}\n"
        f"- Active sessions: {stats['active_sessions']}\n"
        f"- Retention: {stats['retention_hours']} hours"
    )
    return [TextContent(type="text", text=result)]


# ── Server Startup ──

async def run_stdio_server():
    """Run the MCP server using stdio transport."""
    log.info("Starting MCP Office Server (stdio transport)")
    file_cleanup.start()

    async with stdio_server() as (read_stream, write_stream):
        await app.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="mcp-office",
                server_version="1.0.0",
                capabilities=ServerCapabilities(),
            ),
        )

    file_cleanup.stop()
    log.info("MCP Office Server stopped")


async def run_sse_server():
    """Run the MCP server using SSE transport (for remote/server deployment)."""
    import uvicorn
    from starlette.applications import Starlette
    from starlette.routing import Route, Mount
    from starlette.requests import Request

    log.info("Starting MCP Office Server (SSE transport on port 8765)")
    file_cleanup.start()

    sse_transport = SseServerTransport("/messages/")

    # ASGI app for SSE — handles its own HTTP response lifecycle
    async def sse_asgi(scope, receive, send):
        if scope["type"] != "http":
            return
        log.info("=== New SSE connection ===")
        try:
            async with sse_transport.connect_sse(
                scope, receive, send
            ) as (read_stream, write_stream):
                log.info("SSE streams established, starting app.run()")
                await app.run(
                    read_stream,
                    write_stream,
                    InitializationOptions(
                        server_name="mcp-office",
                        server_version="1.0.0",
                        capabilities=ServerCapabilities(),
                    ),
                )
                log.info("app.run() completed normally")
        except Exception as e:
            log.error(f"SSE connection error: {e}")

    # ASGI app for messages — handle_post_message manages its own response via send()
    async def messages_asgi(scope, receive, send):
        if scope["type"] != "http":
            return
        if scope["method"] != "POST":
            from starlette.responses import PlainTextResponse
            response = PlainTextResponse("Method Not Allowed", status_code=405)
            await response(scope, receive, send)
            return
        log.info("Message POST received")
        await sse_transport.handle_post_message(scope, receive, send)

    # File download handler
    async def files_asgi(scope, receive, send):
        if scope["type"] != "http":
            return
        import pathlib
        from starlette.responses import FileResponse, PlainTextResponse
        
        path = scope.get("path", "")
        # Path format: /files/{session_id}/{filename}
        parts = path.split("/")
        if len(parts) < 4 or parts[1] != "files":
            response = PlainTextResponse("Not Found", status_code=404)
            await response(scope, receive, send)
            return
        
        session_id = parts[2]
        filename = "/".join(parts[3:])
        file_path = pathlib.Path(os.environ.get("OUTPUT_DIR", "outputs")) / session_id / filename
        
        if not file_path.is_file():
            response = PlainTextResponse(f"File not found: {filename}", status_code=404)
            await response(scope, receive, send)
            return
        
        log.info(f"File download: {file_path}")
        response = FileResponse(str(file_path), filename=filename)
        await response(scope, receive, send)

    # Single ASGI router — avoids Starlette Mount 307 redirect on /sse → /sse/
    async def router_asgi(scope, receive, send):
        if scope["type"] != "http":
            return
        path = scope.get("path", "")
        if path == "/sse" or path == "/sse/":
            await sse_asgi(scope, receive, send)
        elif path.startswith("/messages"):
            await messages_asgi(scope, receive, send)
        elif path.startswith("/files"):
            await files_asgi(scope, receive, send)
        else:
            from starlette.responses import PlainTextResponse
            response = PlainTextResponse("Not Found", status_code=404)
            await response(scope, receive, send)

    starlette_app = Starlette(
        debug=True,
        routes=[
            Mount("/", app=router_asgi),
        ],
    )

    config = uvicorn.Config(starlette_app, host="0.0.0.0", port=8765, log_level="info")
    server = uvicorn.Server(config)
    await server.serve()

    file_cleanup.stop()
    log.info("MCP Office Server (SSE) stopped")


async def main():
    """Main entry point — supports both stdio and SSE transports."""
    transport = os.environ.get("MCP_TRANSPORT", "stdio").lower()
    if transport == "sse":
        await run_sse_server()
    else:
        await run_stdio_server()


if __name__ == "__main__":
    asyncio.run(main())