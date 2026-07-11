"""Gradient engine — applies linear & radial gradient fills to shapes and slide backgrounds via direct XML manipulation.

Supports:
- Linear gradients (any angle)
- Radial gradients (any center + radius)
- Multi-stop gradients (2+ colors)
- Apply to shapes, text, or slide backgrounds
"""

from __future__ import annotations

from typing import Optional
from pptx.dml.color import RGBColor
from pptx.util import Emu
from pptx.oxml.ns import qn, nsmap
from lxml import etree

from src.utils.colors import hex_to_rgbcolor_tuple


class GradientDef:
    """Definition of a gradient fill."""

    def __init__(
        self,
        gradient_type: str = "linear",
        colors: list[str] = None,
        stops: list[float] = None,
        angle: float = 0.0,
        center_x: float = 0.5,
        center_y: float = 0.5,
        radius: float = 0.5,
    ):
        self.gradient_type = gradient_type  # "linear" or "radial"
        self.colors = colors or ["#1E40AF", "#1E3A8A"]
        self.stops = stops or None
        self.angle = angle
        self.center_x = center_x
        self.center_y = center_y
        self.radius = radius

    @classmethod
    def linear(cls, colors: list[str], angle: float = 0.0, stops: Optional[list[float]] = None) -> GradientDef:
        """Create a linear gradient."""
        return cls(gradient_type="linear", colors=colors, angle=angle, stops=stops)

    @classmethod
    def radial(cls, colors: list[str], center_x: float = 0.5, center_y: float = 0.5, radius: float = 0.5) -> GradientDef:
        """Create a radial gradient."""
        return cls(gradient_type="radial", colors=colors, center_x=center_x, center_y=center_y, radius=radius)

    @classmethod
    def theme_preset(cls, theme_name: str, variant: str = "primary") -> GradientDef:
        """Get a preset gradient for a given theme."""
        presets = {
            "corporate": {
                "primary": ("linear", ["#1E3A8A", "#1E40AF"], 135),
                "accent": ("linear", ["#F59E0B", "#D97706"], 135),
                "dark": ("linear", ["#0F172A", "#1E293B"], 180),
            },
            "dark": {
                "primary": ("linear", ["#1E293B", "#0F172A"], 180),
                "accent": ("linear", ["#3B82F6", "#2563EB"], 135),
                "dark": ("linear", ["#0F172A", "#020617"], 180),
            },
            "creative": {
                "primary": ("linear", ["#78350F", "#92400E"], 135),
                "accent": ("linear", ["#D97706", "#B45309"], 135),
                "warm": ("linear", ["#FFFBEB", "#FEF3C7"], 90),
            },
        }
        theme_presets = presets.get(theme_name, presets["corporate"])
        preset = theme_presets.get(variant, theme_presets["primary"])
        return cls(gradient_type=preset[0], colors=[preset[1][0], preset[1][1]], angle=preset[2])


class GradientEngine:
    """Apply gradient fills to shapes and slide backgrounds."""

    @staticmethod
    def _create_grad_fill_xml(grad: GradientDef) -> etree.Element:
        """Build an a:gradFill XML element from a GradientDef."""
        gf = etree.Element(qn("a:gradFill"))

        if grad.gradient_type == "linear":
            gs_elem = etree.SubElement(gf, qn("a:lin"))
            gs_elem.set("ang", str(int(grad.angle * 60000)))
            gs_elem.set("scaled", "0")
        else:  # radial
            gs_elem = etree.SubElement(gf, qn("a:pathGrad"))
            gs_elem.set("path", "circle")
            fill_rect = etree.SubElement(gs_elem, qn("a:fillToRect"))
            fill_rect.set("l", str(int(grad.center_x * 100000)))
            fill_rect.set("t", str(int(grad.center_y * 100000)))
            fill_rect.set("r", str(int((grad.center_x + grad.radius) * 100000)))
            fill_rect.set("b", str(int((grad.center_y + grad.radius) * 100000)))

        # Create gradient stops
        gs_lst = etree.SubElement(gf, qn("a:gsLst"))
        num_colors = len(grad.colors)
        if num_colors < 2:
            grad.colors = [grad.colors[0], grad.colors[0]] if grad.colors else ["#1E40AF", "#1E3A8A"]
            num_colors = 2

        stops = grad.stops if grad.stops and len(grad.stops) == num_colors else [
            i * 100000 // (num_colors - 1) for i in range(num_colors)
        ]

        for i, (color_hex, pos) in enumerate(zip(grad.colors, stops)):
            gs = etree.SubElement(gs_lst, qn("a:gs"))
            gs.set("pos", str(pos))
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            srgb.set("val", color_hex.lstrip("#"))

        return gf

    @staticmethod
    def apply_to_shape(shape, gradient: GradientDef) -> None:
        """Apply a gradient fill to any shape that has a fill."""
        try:
            sp_pr = shape._element.find(qn("p:spPr"))
            if sp_pr is None:
                sp_pr = etree.SubElement(shape._element, qn("p:spPr"))
            # Remove existing fill
            for fill_elem in sp_pr.findall(qn("a:solidFill")):
                sp_pr.remove(fill_elem)
            for fill_elem in sp_pr.findall(qn("a:gradFill")):
                sp_pr.remove(fill_elem)

            grad_xml = GradientEngine._create_grad_fill_xml(gradient)
            sp_pr.insert(0, grad_xml)
        except Exception:
            pass

    @staticmethod
    def apply_to_slide_bg(slide, gradient: GradientDef) -> None:
        """Apply a gradient fill to the entire slide background."""
        try:
            bg = slide._element.find(qn("p:bg"))
            if bg is None:
                bg = etree.SubElement(slide._element, qn("p:bg"))
            bg_pr = bg.find(qn("p:bgPr"))
            if bg_pr is None:
                bg_pr = etree.SubElement(bg, qn("p:bgPr"))

            for fill_elem in bg_pr.findall(qn("a:solidFill")):
                bg_pr.remove(fill_elem)
            for fill_elem in bg_pr.findall(qn("a:gradFill")):
                bg_pr.remove(fill_elem)

            grad_xml = GradientEngine._create_grad_fill_xml(gradient)
            bg_pr.insert(0, grad_xml)
        except Exception:
            pass

    @staticmethod
    def apply_to_text_frame(text_frame, gradient: GradientDef) -> None:
        """Apply gradient fill to all text in a text frame."""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    r_pr = run._r.find(qn("a:rPr"))
                    if r_pr is None:
                        r_pr = etree.SubElement(run._r, qn("a:rPr"))
                    for fill_elem in r_pr.findall(qn("a:solidFill")):
                        r_pr.remove(fill_elem)
                    for fill_elem in r_pr.findall(qn("a:gradFill")):
                        r_pr.remove(fill_elem)
                    grad_xml = GradientEngine._create_grad_fill_xml(gradient)
                    r_pr.insert(0, grad_xml)
        except Exception:
            pass


# ── Convenience Functions ──

def apply_gradient_to_shape(shape, gradient_type: str = "linear", colors: list = None, angle: float = 135):
    """Quick apply gradient to a shape."""
    grad = GradientDef(gradient_type=gradient_type, colors=colors or ["#1E40AF", "#3B82F6"], angle=angle)
    GradientEngine.apply_to_shape(shape, grad)


def apply_gradient_to_slide(slide, gradient_type: str = "linear", colors: list = None, angle: float = 180):
    """Quick apply gradient to slide background."""
    grad = GradientDef(gradient_type=gradient_type, colors=colors or ["#1E3A8A", "#0F172A"], angle=angle)
    GradientEngine.apply_to_slide_bg(slide, grad)
