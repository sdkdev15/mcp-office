"""NeoShape library — premium shapes beyond python-pptx built-ins.

Supports:
- Pill / Capsule shapes
- Chevron / Arrow process shapes
- Hexagon / Honeycomb
- Callout / Speech bubble
- Star burst
- Diagonal cut corners
- Bracket / Curly brace

All shapes use freeform (a:custGeom) or adjusted preset geometry (a:prstGeom)
via direct XML manipulation.
"""

from __future__ import annotations

from typing import Optional
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from src.visual.gradient_engine import GradientEngine, GradientDef
from src.visual.shadow_engine import ShadowEngine, ShadowDef


class NeoShape:
    """Factory for creating premium shapes on slides."""

    @staticmethod
    def _make_shape(slide, left: float, top: float, width: float, height: float) -> object:
        """Create base shape placeholder."""
        return slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )

    @staticmethod
    def pill(slide, left: float, top: float, width: float, height: float) -> object:
        """Create a pill/capsule shape (rectangle with fully rounded ends).
        Uses MSO_SHAPE.ROUNDED_RECTANGLE with corner ratio set for pill effect.
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        # Adjust corner radius to be half the height for pill effect
        try:
            sp_pr = shape._element.find(qn("p:spPr"))
            if sp_pr is not None:
                prst_geom = sp_pr.find(qn("a:prstGeom"))
                if prst_geom is not None:
                    av_lst = prst_geom.find(qn("a:avLst"))
                    if av_lst is None:
                        av_lst = etree.SubElement(prst_geom, qn("a:avLst"))
                    gd = etree.SubElement(av_lst, qn("a:gd"))
                    gd.set("name", "adj")
                    gd.set("fmla", f"val {int(height * 914400 / 2)}")  # Half height in EMU
        except Exception:
            pass
        return shape

    @staticmethod
    def chevron(slide, left: float, top: float, width: float, height: float, arrow_size: float = 0.3) -> object:
        """Create a chevron/arrow shape for process flows.
        Uses freeform with 6 points: top-left, top-mid, right, bot-mid, bot-left.
        """
        # Use a modified isosceles triangle as base
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        return shape

    @staticmethod
    def hexagon(slide, left: float, top: float, width: float, height: float) -> object:
        """Create a hexagon shape."""
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        except Exception:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        return shape

    @staticmethod
    def callout(slide, left: float, top: float, width: float, height: float) -> object:
        """Create a speech bubble / callout shape."""
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        except Exception:
            shape = NeoShape.pill(slide, left, top, width, height)
        return shape

    @staticmethod
    def star(slide, left: float, top: float, width: float, height: float, points: int = 5) -> object:
        """Create a star shape for highlights."""
        shape_map = {
            4: MSO_SHAPE.STAR_4_POINT,
            5: MSO_SHAPE.STAR_5_POINT,
            8: MSO_SHAPE.STAR_8_POINT,
        }
        mso = shape_map.get(points, MSO_SHAPE.STAR_5_POINT)
        try:
            shape = slide.shapes.add_shape(
                mso,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        except Exception:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        return shape

    @staticmethod
    def bracket(slide, left: float, top: float, width: float, height: float) -> object:
        """Create a bracket/curly brace shape for quotes."""
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.LEFT_BRACKET,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        except Exception:
            shape = NeoShape.pill(slide, left, top, 0.08, height)
        return shape

    @staticmethod
    def diagonal_cut(slide, left: float, top: float, width: float, height: float, corner: str = "tr") -> object:
        """Create a rectangle with a diagonal corner cut.
        corner options: 'tr' (top-right), 'tl', 'br', 'bl'
        """
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.SNIP_1_CORNER_RECTANGLE if corner in ("tr", "tl")
                else MSO_SHAPE.SNIP_2_CORNER_RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        except Exception:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        return shape

    @staticmethod
    def circle(slide, left: float, top: float, size: float) -> object:
        """Create a perfect circle."""
        return slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top), Inches(size), Inches(size),
        )

    @staticmethod
    def line(slide, x1: float, y1: float, x2: float, y2: float, color: str = "#CBD5E1", width: float = 1.0) -> object:
        """Create a connector line with color."""
        connector = slide.shapes.add_connector(
            MSO_SHAPE.STRAIGHT_CONNECTOR,  # type: ignore
            Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        try:
            connector.line.color.rgb = RGBColor(
                *__import__('src.utils.colors', fromlist=['hex_to_rgbcolor_tuple']).hex_to_rgbcolor_tuple(color)
            )
            connector.line.width = Pt(width)
        except Exception:
            pass
        return connector

    @staticmethod
    def apply_style(shape, fill_color: str = None, border_color: str = None, border_width: float = 0,
                    gradient: Optional[GradientDef] = None, shadow: Optional[ShadowDef] = None) -> None:
        """Apply styling to a shape: fill, border, gradient, shadow."""
        if gradient:
            GradientEngine.apply_to_shape(shape, gradient)
        elif fill_color:
            shape.fill.solid()
            from pptx.dml.color import RGBColor
            from src.utils.colors import hex_to_rgbcolor_tuple
            r, g, b = hex_to_rgbcolor_tuple(fill_color)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)

        if border_color:
            shape.line.color.rgb = RGBColor(
                *__import__('src.utils.colors', fromlist=['hex_to_rgbcolor_tuple']).hex_to_rgbcolor_tuple(border_color)
            )
            shape.line.width = Pt(border_width)
        elif border_width == 0:
            shape.line.fill.background()

        if shadow:
            ShadowEngine.apply_to_shape(shape, shadow)


# ── Convenience Functions ──

def add_pill_card(slide, left: float, top: float, width: float, height: float,
                  fill: str = "#1E40AF", shadow: bool = True):
    """Quick pill-shaped card with optional shadow."""
    shape = NeoShape.pill(slide, left, top, width, height)
    if fill:
        shape.fill.solid()
        from pptx.dml.color import RGBColor
        from src.utils.colors import hex_to_rgbcolor_tuple
        r, g, b = hex_to_rgbcolor_tuple(fill)
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()
    if shadow:
        ShadowEngine.apply_to_shape(shape, ShadowDef.soft())
    return shape


def add_numbered_circle(slide, number: int, left: float, top: float, size: float = 0.5,
                        bg_color: str = "#1E40AF", text_color: str = "#FFFFFF", font_size: int = 14):
    """Quick numbered circle for agenda items."""
    shape = NeoShape.circle(slide, left, top, size)
    shape.fill.solid()
    from pptx.dml.color import RGBColor
    from src.utils.colors import hex_to_rgbcolor_tuple
    r, g, b = hex_to_rgbcolor_tuple(bg_color)
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*hex_to_rgbcolor_tuple(text_color))
    p.alignment = 1  # center
    return shape
