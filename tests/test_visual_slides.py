"""Tests for visual slide builder, visual elements, and auto-icon features."""

import os
import pytest
from pathlib import Path
from pptx import Presentation

from src.generators.visual_slide_builder import VisualSlideBuilder
from src.generators.pptx_generator import PPTXGenerator
from src.generators.docx_generator import DOCXGenerator
from src.utils.visual_elements import VisualElements
from src.utils.icon_library import find_icon_by_keyword


class TestVisualSlideBuilder:
    """Test premium visual slide creation."""

    @pytest.fixture
    def builder(self):
        return VisualSlideBuilder(theme_name="corporate")

    def _new_pres(self):
        return Presentation()

    def test_build_cover_slide(self, builder):
        """build_cover_slide should create a slide without errors."""
        pres = self._new_pres()
        builder.build_cover_slide(
            pres=pres,
            title="Q4 Security Report",
            subtitle="Incident Analysis & Response",
            icon_name="shield",
        )
        assert len(pres.slides) == 1

    def test_build_cover_auto_icon(self, builder):
        """build_cover_slide should auto-detect icon from title."""
        pres = self._new_pres()
        builder.build_cover_slide(
            pres=pres,
            title="Database Migration Plan",
            subtitle="Phase 1",
        )
        assert len(pres.slides) == 1

    def test_build_agenda_slide(self, builder):
        """build_agenda_slide should create agenda layout."""
        pres = self._new_pres()
        items = [
            {"title": "Introduction", "subtitle": "Overview"},
            {"title": "Analysis", "subtitle": "Deep dive"},
            {"title": "Conclusion", "subtitle": "Wrap-up"},
        ]
        builder.build_agenda_slide(pres=pres, title="Agenda", items=items)
        assert len(pres.slides) == 1

    def test_build_exec_summary(self, builder):
        """build_exec_summary_slide should create stat boxes + body."""
        pres = self._new_pres()
        boxes = [
            {"number": "99.9%", "label": "Uptime", "sub_label": "Last 30 days"},
            {"number": "12", "label": "Incidents", "sub_label": "Resolved"},
            {"number": "4.2s", "label": "Avg Response", "sub_label": "API p99"},
        ]
        builder.build_exec_summary_slide(
            pres=pres,
            title="Executive Summary",
            stat_boxes=boxes,
            body_text="All metrics within SLA targets.",
        )
        assert len(pres.slides) == 1

    def test_build_timeline_slide(self, builder):
        """build_timeline_slide should create timeline layout."""
        pres = self._new_pres()
        events = [
            {"title": "Planning", "date": "2026-01"},
            {"title": "Development", "date": "2026-03"},
            {"title": "Launch", "date": "2026-06"},
        ]
        builder.build_timeline_slide(
            pres=pres,
            title="Project Timeline",
            events=events,
            horizontal=True,
        )
        assert len(pres.slides) == 1

    def test_build_flow_slide(self, builder):
        """build_flow_slide should create flow diagram."""
        pres = self._new_pres()
        nodes = [
            {"label": "Client", "icon": "globe"},
            {"label": "API Gateway", "icon": "api"},
            {"label": "Backend", "icon": "server"},
        ]
        connections = [(0, 1), (1, 2)]
        builder.build_flow_slide(
            pres=pres,
            title="Architecture",
            nodes=nodes,
            connections=connections,
        )
        assert len(pres.slides) == 1

    def test_visual_slide_saves(self, builder, tmp_path):
        """A presentation with visual slides should save to file."""
        pres = self._new_pres()
        builder.build_cover_slide(pres=pres, title="Test Deck", subtitle="Auto-generated")
        builder.build_agenda_slide(pres=pres, title="Agenda", items=[{"title": "Topic 1"}])

        out = tmp_path / "test.pptx"
        pres.save(str(out))
        assert out.exists()
        assert out.stat().st_size > 0

    def test_build_exec_summary_empty(self, builder):
        """build_exec_summary_slide with no stat boxes should handle gracefully."""
        pres = self._new_pres()
        builder.build_exec_summary_slide(pres=pres, title="Empty", stat_boxes=[], body_text="Nothing to show.")
        assert len(pres.slides) == 1


class TestVisualElements:
    """Test visual element generation."""

    @pytest.fixture
    def visuals(self, tmp_path):
        return VisualElements(cache_dir=str(tmp_path))

    def test_stat_box(self, visuals):
        """stat_box should render a PNG."""
        path = visuals.stat_box(number="42", label="Total Issues", bg_color="#1E40AF")
        assert os.path.exists(path)
        assert path.endswith(".png")

    def test_stat_boxes_row(self, visuals):
        """stat_boxes_row should render a row of stat boxes."""
        boxes = [
            {"number": "100", "label": "Uptime %"},
            {"number": "0", "label": "Outages"},
        ]
        path = visuals.stat_boxes_row(boxes=boxes)
        assert os.path.exists(path)

    def test_agenda_card(self, visuals):
        """agenda_card should render a card PNG."""
        path = visuals.agenda_card(number=1, title="Introduction", subtitle="Overview and context")
        assert os.path.exists(path)

    def test_timeline_item(self, visuals):
        """timeline_item should render a timeline item."""
        path = visuals.timeline_item(step=1, title="Phase 1", date="2026-01")
        assert os.path.exists(path)

    def test_flow_diagram(self, visuals):
        """flow_diagram should render a flow diagram."""
        nodes = [
            {"label": "A"},
            {"label": "B"},
            {"label": "C"},
        ]
        path = visuals.flow_diagram(nodes=nodes, connections=[(0, 1), (1, 2)])
        assert os.path.exists(path)

    def test_comparison_card(self, visuals):
        """comparison_card should render side-by-side comparison."""
        path = visuals.comparison_card(
            title_left="Option A",
            title_right="Option B",
            items_left=["Fast", "Simple"],
            items_right=["Slow", "Complex"],
        )
        assert os.path.exists(path)

    def test_cover_slide_bg(self, visuals):
        """cover_slide_bg should render a cover slide background."""
        path = visuals.cover_slide_bg()
        assert os.path.exists(path)

    def test_numbered_circle(self, visuals):
        """numbered_circle should render a numbered circle."""
        path = visuals.numbered_circle(number=5, size=48)
        assert os.path.exists(path)


class TestPPTXAutoIcons:
    """Test auto-icons end-to-end with PPTXGenerator."""

    def test_pptx_with_auto_icons(self, tmp_path):
        """PPTXGenerator should create slides with auto_icons enabled."""
        gen = PPTXGenerator(theme_name="corporate")
        data = gen.create_presentation(
            title="Test",
            slides=[
                {"title": "Security Overview", "content": "Lock and shield analysis", "auto_icon": True},
                {"title": "Database Stats", "content": "Storage metrics", "auto_icon": True},
            ],
            auto_icons=True,
        )
        assert isinstance(data, bytes)
        assert len(data) > 0

    def test_pptx_visual_slides(self, tmp_path):
        """PPTXGenerator should create visual slides via add_slide."""
        gen = PPTXGenerator(theme_name="corporate")
        data = gen.create_presentation(
            title="Visual Test",
            slides=[
                {"title": "Cover Title", "visual_type": "cover", "content": "Subtitle", "icon_name": "shield"},
                {"title": "Agenda", "visual_type": "agenda", "agenda_items": [{"title": "Topic 1"}]},
                {
                    "title": "Summary",
                    "visual_type": "exec-summary",
                    "stat_boxes": [{"number": "10", "label": "Items"}],
                    "content": "Body text",
                },
                {
                    "title": "Timeline",
                    "visual_type": "timeline",
                    "timeline_events": [{"title": "Event 1", "date": "2026-01"}],
                },
                {
                    "title": "Architecture",
                    "visual_type": "flow",
                    "flow_nodes": [{"label": "Node A"}, {"label": "Node B"}],
                    "flow_connections": [[0, 1]],
                },
            ],
        )
        assert isinstance(data, bytes)
        assert len(data) > 0


class TestDOCXAutoIcons:
    """Test auto-icons end-to-end with DOCXGenerator."""

    def test_docx_with_auto_icons(self, tmp_path):
        """DOCXGenerator should create sections with auto_icons."""
        gen = DOCXGenerator(theme_name="corporate")
        data = gen.create_document_with_content(
            title="Icon Test",
            sections=[
                {"type": "heading_1", "text": "Security Analysis", "icon": "shield"},
                {"type": "icon_heading", "text": "Server Overview", "icon": "server", "level": 2},
                {"type": "paragraph", "text": "This document uses icons."},
            ],
            auto_icons=True,
        )
        assert isinstance(data, bytes)
        assert len(data) > 0

    def test_docx_cover_page(self, tmp_path):
        """DOCXGenerator should create a cover page section."""
        gen = DOCXGenerator(theme_name="corporate")
        data = gen.create_document_with_content(
            title="Cover Test",
            sections=[
                {
                    "type": "cover_page",
                    "text": "Quarterly Report",
                    "subtitle": "Q3 2026",
                    "icon": "chart-bar",
                },
                {"type": "heading_1", "text": "Introduction"},
                {"type": "paragraph", "text": "Content goes here."},
            ],
        )
        assert isinstance(data, bytes)
        assert len(data) > 0

    def test_docx_summary_card(self, tmp_path):
        """DOCXGenerator should create summary card sections."""
        gen = DOCXGenerator(theme_name="corporate")
        data = gen.create_document_with_content(
            title="Summary Test",
            sections=[
                {
                    "type": "summary_card",
                    "text": "Key Findings",
                    "icon": "alert",
                    "items": ["Finding 1", "Finding 2", "Finding 3"],
                },
                {"type": "heading_2", "text": "Details"},
                {"type": "list_bullet", "items": ["Detail A", "Detail B"]},
            ],
        )
        assert isinstance(data, bytes)
        assert len(data) > 0

    def test_docx_full_sections(self, tmp_path):
        """DOCXGenerator should handle all section types together."""
        gen = DOCXGenerator(theme_name="corporate")
        sections = [
            {"type": "title", "text": "Full Document Test"},
            {"type": "subtitle", "text": "All section types"},
            {"type": "toc"},
            {"type": "heading_1", "text": "Infrastructure"},
            {"type": "paragraph", "text": "Server and network overview."},
            {"type": "icon_heading", "text": "Database Performance", "icon": "database", "level": 2},
            {
                "type": "table",
                "headers": ["Metric", "Value"],
                "rows": [["Latency", "10ms"], ["Throughput", "5000 rps"]],
            },
            {"type": "list_bullet", "items": ["Item 1", "Item 2"]},
            {"type": "list_number", "items": ["Step 1", "Step 2"]},
            {
                "type": "summary_card",
                "text": "Summary",
                "icon": "check",
                "items": ["All OK"],
            },
        ]
        data = gen.create_document_with_content(title="Full Test", sections=sections, auto_icons=True)
        assert isinstance(data, bytes)
        assert len(data) > 0